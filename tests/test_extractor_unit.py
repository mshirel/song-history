"""Unit tests for worship_catalog.extractor internal functions."""

import io
import logging
import zipfile
from pathlib import Path

import pytest

from worship_catalog.extractor import (
    _NO_TEXT_STREAK_THRESHOLD,
    _create_song_occurrence,
    _extract_title_candidates,
    _group_song_slides,
    _is_song_title_slide,
    extract_songs,
)
from worship_catalog.normalize import _TITLE_MAX_LENGTH
from worship_catalog.pptx_reader import Slide, SlideImage, SlideText


def make_slide(
    index: int = 0,
    lines: list[str] | None = None,
    hidden: bool = False,
    image_blobs: list[bytes | None] | None = None,
) -> Slide:
    text = SlideText(text_lines=lines or [])
    images = [SlideImage(shape_id=i + 1, blob=b) for i, b in enumerate(image_blobs or [])]
    return Slide(index=index, hidden=hidden, text=text, images=images)


class TestExtractTitleCandidates:
    def test_empty_slide_returns_empty(self):
        slide = make_slide(lines=[])
        assert _extract_title_candidates(slide) == []

    def test_copyright_line_filtered(self):
        slide = make_slide(lines=["Copyright 2020 by Acme Publishing"])
        assert _extract_title_candidates(slide) == []

    def test_long_line_over_120_chars_filtered(self):
        slide = make_slide(lines=["A" * 121])
        assert _extract_title_candidates(slide) == []

    def test_single_char_filtered(self):
        slide = make_slide(lines=["A"])
        assert _extract_title_candidates(slide) == []

    def test_scripture_reference_filtered(self):
        slide = make_slide(lines=["John 3:16"])
        assert _extract_title_candidates(slide) == []

    def test_normal_title_returned(self):
        slide = make_slide(lines=["Amazing Grace"])
        assert "Amazing Grace" in _extract_title_candidates(slide)

    def test_offering_line_filtered(self):
        slide = make_slide(lines=["Giving Online"])
        assert _extract_title_candidates(slide) == []

    def test_multiple_lines_mixed(self):
        slide = make_slide(lines=["Amazing Grace", "Copyright 2020", "How Sweet the Sound"])
        candidates = _extract_title_candidates(slide)
        assert "Amazing Grace" in candidates
        assert "How Sweet the Sound" in candidates
        assert len([c for c in candidates if "copyright" in c.lower()]) == 0


class TestIsSongTitleSlide:
    def test_paperlesshymnal_marker_returns_true(self):
        slide = make_slide(lines=["Amazing Grace", "PaperlessHymnal.com"])
        assert _is_song_title_slide(slide) is True

    def test_taylor_publications_marker_returns_true(self):
        slide = make_slide(lines=["How Great Thou Art", "Taylor Publications"])
        assert _is_song_title_slide(slide) is True

    def test_verse_prefix_returns_true(self):
        slide = make_slide(lines=["1 – Amazing grace, how sweet the sound"])
        assert _is_song_title_slide(slide) is True

    def test_chorus_prefix_returns_true(self):
        slide = make_slide(lines=["Chorus – How great thou art"])
        assert _is_song_title_slide(slide) is True

    def test_plain_prose_returns_false(self):
        slide = make_slide(lines=["This is a sermon slide about faith and hope."])
        assert _is_song_title_slide(slide) is False

    def test_empty_slide_returns_false(self):
        slide = make_slide(lines=[])
        assert _is_song_title_slide(slide) is False


class TestGroupSongSlides:
    def test_empty_list_returns_empty(self):
        assert _group_song_slides([]) == []

    def test_single_titled_song_multiple_slides(self):
        """Multiple slides with same canonical title stay in one group."""
        slides = [
            make_slide(0, ["Amazing Grace", "PaperlessHymnal.com"]),
            make_slide(1, ["Amazing Grace"]),  # same canonical — continues group
            make_slide(2, ["Amazing Grace"]),
        ]
        groups = _group_song_slides(slides)
        assert len(groups) == 1
        assert groups[0][0] == "amazing grace"

    def test_two_songs_with_title_slides(self):
        slides = [
            make_slide(0, ["Amazing Grace", "PaperlessHymnal.com"]),
            make_slide(1, ["Amazing Grace"]),
            make_slide(2, ["How Great Thou Art", "PaperlessHymnal.com"]),
            make_slide(3, ["How Great Thou Art"]),
        ]
        groups = _group_song_slides(slides)
        assert len(groups) == 2
        canonicals = [c for c, _ in groups]
        assert "amazing grace" in canonicals
        assert "how great thou art" in canonicals

    def test_five_consecutive_empty_slides_closes_group(self):
        slides = [
            make_slide(0, ["Amazing Grace", "PaperlessHymnal.com"]),
            make_slide(1, []),  # image-only
            make_slide(2, []),
            make_slide(3, []),
            make_slide(4, []),
            make_slide(5, []),  # 5th empty — closes group
            make_slide(6, ["How Great Thou Art", "PaperlessHymnal.com"]),
        ]
        groups = _group_song_slides(slides)
        assert len(groups) == 2

    def test_skip_offering_slide(self):
        slides = [
            make_slide(0, ["Giving Online"]),
        ]
        groups = _group_song_slides(slides)
        assert groups == []


class TestCreateSongOccurrence:
    def test_returns_song_occurrence_with_correct_ordinal(self):
        slides = [make_slide(0, ["Amazing Grace", "PaperlessHymnal.com"])]
        result = _create_song_occurrence(1, "amazing grace", slides)
        assert result.ordinal == 1

    def test_display_title_extracted_from_slide(self):
        slides = [make_slide(0, ["Amazing Grace", "PaperlessHymnal.com"])]
        result = _create_song_occurrence(1, "amazing grace", slides)
        assert result.display_title == "Amazing Grace"

    def test_falls_back_to_canonical_when_no_text(self):
        slides = [make_slide(0, [])]
        result = _create_song_occurrence(1, "amazing grace", slides)
        assert result.display_title == "amazing grace"

    def test_slide_range_set_correctly(self):
        slides = [
            make_slide(3, ["1 – Amazing grace"]),
            make_slide(4, ["2 – That saved a wretch"]),
        ]
        result = _create_song_occurrence(1, "amazing grace", slides)
        assert result.first_slide_index == 3
        assert result.last_slide_index == 4

    def test_ocr_not_called_when_credits_found_in_text(self, monkeypatch):
        """If credits are in text, OCR should not be called."""
        called = []
        monkeypatch.setattr(
            "worship_catalog.extractor.extract_credits_via_vision",
            lambda blob: called.append(blob) or "Words: X",
        )
        slides = [
            make_slide(0, ["Amazing Grace", "Words: John Newton", "PaperlessHymnal.com"],
                       image_blobs=[b"\xff\xd8"])
        ]
        _create_song_occurrence(1, "amazing grace", slides, use_ocr=True)
        assert called == [], "OCR should not be called when credits found in text"

    def test_ocr_not_called_when_no_image(self, monkeypatch):
        """If slide has no image, OCR should not be called even when use_ocr=True."""
        called = []
        monkeypatch.setattr(
            "worship_catalog.extractor.extract_credits_via_vision",
            lambda blob: called.append(blob) or "Words: X",
        )
        slides = [make_slide(0, ["Amazing Grace"])]
        _create_song_occurrence(1, "amazing grace", slides, use_ocr=True)
        assert called == []


class TestExtractSongsFileHash:
    """Tests for ExtractionResult.file_hash — issue #15."""

    def _minimal_pptx(self) -> bytes:
        """Build a minimal valid PPTX (ZIP) with a single blank slide."""
        ns_pkg = "http://schemas.openxmlformats.org/package/2006"
        ns_doc = "http://schemas.openxmlformats.org/officeDocument/2006"
        ns_pml = "http://schemas.openxmlformats.org/presentationml/2006/main"
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
            zf.writestr(
                "[Content_Types].xml",
                '<?xml version="1.0" encoding="UTF-8"?>'
                f'<Types xmlns="{ns_pkg}/content-types">'
                '<Default Extension="rels" ContentType='
                f'"{ns_pkg}/relationships+xml"/>'
                '<Default Extension="xml" ContentType="application/xml"/>'
                '<Override PartName="/ppt/presentation.xml" ContentType='
                '"application/vnd.openxmlformats-officedocument'
                '.presentationml.presentation.main+xml"/>'
                "</Types>",
            )
            zf.writestr(
                "_rels/.rels",
                '<?xml version="1.0" encoding="UTF-8"?>'
                f'<Relationships xmlns="{ns_pkg}/relationships">'
                f'<Relationship Id="rId1" Type="{ns_doc}/relationships'
                '/officeDocument" Target="ppt/presentation.xml"/>'
                "</Relationships>",
            )
            zf.writestr(
                "ppt/presentation.xml",
                '<?xml version="1.0" encoding="UTF-8"?>'
                f'<p:presentation xmlns:p="{ns_pml}"'
                f' xmlns:r="{ns_doc}/relationships">'
                '<p:sldMasterIdLst/>'
                '<p:sldSz cx="9144000" cy="5143500"/>'
                '<p:notesSz cx="6858000" cy="9144000"/>'
                "</p:presentation>",
            )
            zf.writestr(
                "ppt/_rels/presentation.xml.rels",
                '<?xml version="1.0" encoding="UTF-8"?>'
                f'<Relationships xmlns="{ns_pkg}/relationships">'
                "</Relationships>",
            )
        return buf.getvalue()

    def test_file_hash_is_non_empty(self, tmp_path: Path) -> None:
        """ExtractionResult.file_hash must be a non-empty string."""
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(self._minimal_pptx())
        result = extract_songs(pptx_path)
        assert result.file_hash, "file_hash must not be empty"
        assert isinstance(result.file_hash, str)

    def test_file_hash_is_stable(self, tmp_path: Path) -> None:
        """Same file produces the same hash on repeated calls."""
        pptx_bytes = self._minimal_pptx()
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(pptx_bytes)
        result1 = extract_songs(pptx_path)
        result2 = extract_songs(pptx_path)
        assert result1.file_hash == result2.file_hash

    def test_different_files_produce_different_hashes(self, tmp_path: Path) -> None:
        """Different file content produces different hashes."""
        bytes_a = self._minimal_pptx()
        bytes_b = bytes_a + b"\x00"  # trivially different content
        path_a = tmp_path / "a.pptx"
        path_b = tmp_path / "b.pptx"
        path_a.write_bytes(bytes_a)
        path_b.write_bytes(bytes_b)
        result_a = extract_songs(path_a)
        result_b = extract_songs(path_b)
        assert result_a.file_hash != result_b.file_hash


class TestPptxSizeLimits:
    """Tests for PPTX pre-flight size and slide-count checks — issue #40."""

    def test_file_larger_than_limit_is_rejected(self, tmp_path: Path) -> None:
        """Files above MAX_PPTX_SIZE_BYTES are rejected before loading."""
        large_file = tmp_path / "big.pptx"
        # 60 MB — over any sane limit
        large_file.write_bytes(b"PK" + b"\x00" * (60 * 1024 * 1024))
        with pytest.raises(ValueError, match="(?i)exceeds maximum"):
            extract_songs(large_file)

    def test_file_within_limit_is_accepted(self, tmp_path: Path) -> None:
        """A small PPTX within size limit is not rejected by size check."""
        from tests.test_extractor_unit import TestExtractSongsFileHash
        pptx_bytes = TestExtractSongsFileHash()._minimal_pptx()
        small_file = tmp_path / "small.pptx"
        small_file.write_bytes(pptx_bytes)
        # Should not raise a ValueError about size
        try:
            extract_songs(small_file)
        except ValueError as exc:
            assert "exceeds maximum" not in str(exc), f"Unexpected size error: {exc}"

    def test_presentation_with_too_many_slides_is_rejected(
        self, tmp_path: Path
    ) -> None:
        """Presentations with more than MAX_SLIDES slides are rejected."""
        # Build a minimal PPTX then monkey-patch load_pptx to return a mock
        from unittest.mock import MagicMock, patch

        mock_prs = MagicMock()
        # Simulate 1001 slides
        mock_prs.slides = [MagicMock()] * 1001

        pptx_path = tmp_path / "big_slide_count.pptx"
        pptx_path.write_bytes(b"PK\x03\x04")  # minimal ZIP magic

        with patch("worship_catalog.extractor.load_pptx", return_value=mock_prs):
            with pytest.raises(ValueError, match="(?i)too many slides"):
                extract_songs(pptx_path)


class TestNamedConstants:
    """Verify magic numbers have been replaced with named constants (#22)."""

    def test_title_max_length_constant_exists_and_is_120(self):
        assert _TITLE_MAX_LENGTH == 120

    def test_no_text_streak_threshold_exists_and_is_5(self):
        assert _NO_TEXT_STREAK_THRESHOLD == 5

    def test_title_max_length_enforced_in_extractor(self):
        """A line exactly at the limit is still a valid title candidate."""
        slide = make_slide(lines=["A" * _TITLE_MAX_LENGTH])
        candidates = _extract_title_candidates(slide)
        assert len(candidates) > 0

    def test_line_over_title_max_length_excluded(self):
        """A line longer than the limit is rejected as a title candidate."""
        slide = make_slide(lines=["A" * (_TITLE_MAX_LENGTH + 1)])
        candidates = _extract_title_candidates(slide)
        assert len(candidates) == 0


class TestExtractionLogging:
    """Verify extraction decisions are logged for debugging silent data loss (#23)."""

    def test_extract_songs_logs_song_count(self, tmp_path, caplog):
        """extract_songs() should log how many songs were found."""
        import logging
        from tests.test_extractor_unit import TestExtractSongsFileHash

        pptx_bytes = TestExtractSongsFileHash()._minimal_pptx()
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(pptx_bytes)

        with caplog.at_level(logging.DEBUG, logger="worship_catalog.extractor"):
            extract_songs(pptx_path)

        assert any("song" in record.message.lower() for record in caplog.records)

    def test_no_credits_found_is_logged(self, caplog):
        """When no credits are found, a DEBUG message should mention it."""
        import logging
        from worship_catalog.extractor import _create_song_occurrence
        from worship_catalog.pptx_reader import Slide, SlideText

        slide = make_slide(lines=["Amazing Grace"])

        with caplog.at_level(logging.DEBUG, logger="worship_catalog.extractor"):
            _create_song_occurrence(
                ordinal=1,
                canonical_title="amazing grace",
                slides=[slide],
                use_ocr=False,
            )

        messages = [r.message.lower() for r in caplog.records]
        assert any("credit" in m or "no credit" in m for m in messages)

    def test_credits_found_via_text_is_logged(self, caplog):
        """When credits are found via text parsing, it should be logged at DEBUG."""
        import logging
        from worship_catalog.extractor import _create_song_occurrence

        slide = make_slide(lines=["Amazing Grace", "Words: John Newton / Music: Traditional"])

        with caplog.at_level(logging.DEBUG, logger="worship_catalog.extractor"):
            _create_song_occurrence(
                ordinal=1,
                canonical_title="amazing grace",
                slides=[slide],
                use_ocr=False,
            )

        messages = [r.message.lower() for r in caplog.records]
        assert any("credit" in m for m in messages)


# ---------------------------------------------------------------------------
# _try_ocr_credits exception handling (#55)
# ---------------------------------------------------------------------------

class TestTryOcrCreditsExceptionHandling:
    """_try_ocr_credits() must log on failure, not silently swallow errors (#55)."""

    def _slides_with_image(self):
        return [make_slide(0, ["Amazing Grace"], image_blobs=[b"\xff\xd8\xff"])]

    def test_os_error_returns_none_and_logs_warning(self, monkeypatch, caplog):
        import logging
        from worship_catalog.extractor import _try_ocr_credits

        def _raise(blob):
            raise OSError("ANTHROPIC_API_KEY is not set")

        monkeypatch.setattr("worship_catalog.extractor.extract_credits_via_vision", _raise)

        with caplog.at_level(logging.WARNING, logger="worship_catalog.extractor"):
            result = _try_ocr_credits(self._slides_with_image())

        assert result is None
        assert any(r.levelno >= logging.WARNING for r in caplog.records)

    def test_import_error_returns_none_and_logs_warning(self, monkeypatch, caplog):
        import logging
        from worship_catalog.extractor import _try_ocr_credits

        def _raise(blob):
            raise ImportError("anthropic not installed")

        monkeypatch.setattr("worship_catalog.extractor.extract_credits_via_vision", _raise)

        with caplog.at_level(logging.WARNING, logger="worship_catalog.extractor"):
            result = _try_ocr_credits(self._slides_with_image())

        assert result is None
        assert any(r.levelno >= logging.WARNING for r in caplog.records)

    def test_api_error_returns_none_and_logs_warning(self, monkeypatch, caplog):
        """Any other exception (e.g. SDK API error) must be logged, not silently swallowed."""
        import logging
        from worship_catalog.extractor import _try_ocr_credits

        def _raise(blob):
            raise RuntimeError("connection refused")

        monkeypatch.setattr("worship_catalog.extractor.extract_credits_via_vision", _raise)

        with caplog.at_level(logging.WARNING, logger="worship_catalog.extractor"):
            result = _try_ocr_credits(self._slides_with_image())

        assert result is None
        assert any(r.levelno >= logging.WARNING for r in caplog.records)

    def test_error_message_includes_exception_type(self, monkeypatch, caplog):
        """Log message should include the exception type so it's diagnosable."""
        import logging
        from worship_catalog.extractor import _try_ocr_credits

        def _raise(blob):
            raise ValueError("unexpected response shape")

        monkeypatch.setattr("worship_catalog.extractor.extract_credits_via_vision", _raise)

        with caplog.at_level(logging.WARNING, logger="worship_catalog.extractor"):
            _try_ocr_credits(self._slides_with_image())

        messages = " ".join(r.message for r in caplog.records)
        assert "ValueError" in messages or "valueerror" in messages.lower()


# ---------------------------------------------------------------------------
# Timeout tests (#73)
# ---------------------------------------------------------------------------

class TestExtractionTimeout:
    """extract_songs() must raise TimeoutError if it runs longer than _MAX_EXTRACT_SECONDS (#73)."""

    def test_extraction_raises_if_time_limit_exceeded(self, tmp_path, monkeypatch):
        """extract_songs() raises TimeoutError if it runs longer than _MAX_EXTRACT_SECONDS."""
        import worship_catalog.extractor as ext_module
        monkeypatch.setattr(ext_module, "_MAX_EXTRACT_SECONDS", 0.001)

        from pptx import Presentation
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        for _ in range(50):
            prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(pptx_path)

        with pytest.raises(TimeoutError):
            ext_module.extract_songs(pptx_path)

    def test_normal_extraction_completes_within_limit(self, tmp_path):
        """A minimal PPTX completes well within the default time limit."""
        from pptx import Presentation
        from worship_catalog.extractor import extract_songs
        pptx_path = tmp_path / "quick.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(pptx_path)
        # Should not raise — returns ExtractionResult normally
        result = extract_songs(pptx_path)
        assert result is not None


# ---------------------------------------------------------------------------
# CreditResolver abstraction (#53)
# ---------------------------------------------------------------------------


class TestCreditResolver:
    """CreditResolver encapsulates the three-step credit cascade (#53)."""

    def _make_slides(self, image_blob: bytes | None = None) -> list:
        images = [image_blob] if image_blob is not None else []
        return [make_slide(index=0, lines=["Amazing Grace"], image_blobs=images)]

    def test_returns_parsed_credits_when_complete(self):
        """If parsed credits already contain words_by, the resolver returns them immediately."""
        from worship_catalog.extractor import CreditResolver, OcrBudget

        resolver = CreditResolver(library_index=None, ocr_budget=None)
        # Pre-parsed credits that are already complete
        complete = {"words_by": "John Newton", "music_by": None, "arranger": None}
        slides = self._make_slides()
        result = resolver.resolve(slides, complete, canonical_title="amazing grace")
        assert result["words_by"] == "John Newton"

    def test_falls_through_to_library_when_parsed_incomplete(self):
        """When parsed credits are empty, the resolver tries the library index."""
        from worship_catalog.extractor import CreditResolver, OcrBudget

        library = {"amazing grace": {"words_by": "John Newton", "music_by": None, "arranger": None}}
        resolver = CreditResolver(library_index=library, ocr_budget=None)
        empty_credits: dict = {"words_by": None, "music_by": None, "arranger": None}
        slides = self._make_slides()
        result = resolver.resolve(slides, empty_credits, canonical_title="amazing grace")
        assert result["words_by"] == "John Newton"

    def test_library_miss_leaves_credits_empty_when_no_ocr(self):
        """When library misses and no OCR is configured, credits remain empty."""
        from worship_catalog.extractor import CreditResolver

        resolver = CreditResolver(library_index={}, ocr_budget=None)
        empty_credits: dict = {"words_by": None, "music_by": None, "arranger": None}
        slides = self._make_slides()
        result = resolver.resolve(slides, empty_credits, canonical_title="unknown song")
        assert result["words_by"] is None
        assert result["music_by"] is None

    def test_falls_through_to_ocr_when_library_misses(self, monkeypatch):
        """When library has no match, OCR is attempted if budget allows."""
        from worship_catalog.extractor import CreditResolver, OcrBudget
        import worship_catalog.extractor as ext_module

        ocr_calls: list[bytes] = []

        def fake_try_ocr(slides):
            ocr_calls.append(b"called")
            return "Words: Anne Steele"

        monkeypatch.setattr(ext_module, "_try_ocr_credits", fake_try_ocr)

        budget = OcrBudget(max_calls=5)
        resolver = CreditResolver(library_index={}, ocr_budget=budget, use_ocr=True)
        empty_credits: dict = {"words_by": None, "music_by": None, "arranger": None}
        slides = self._make_slides(image_blob=b"\x89PNG\r\n")
        result = resolver.resolve(slides, empty_credits, canonical_title="blest are the pure")
        assert len(ocr_calls) == 1
        assert result["words_by"] == "Anne Steele"

    def test_ocr_not_called_when_budget_exhausted(self, monkeypatch):
        """OCR must not be called when the OcrBudget is exhausted."""
        from worship_catalog.extractor import CreditResolver, OcrBudget
        import worship_catalog.extractor as ext_module

        ocr_calls: list[bytes] = []

        def fake_try_ocr(slides):
            ocr_calls.append(b"called")
            return "Words: Somebody"

        monkeypatch.setattr(ext_module, "_try_ocr_credits", fake_try_ocr)

        budget = OcrBudget(max_calls=0)  # budget already exhausted
        resolver = CreditResolver(library_index={}, ocr_budget=budget, use_ocr=True)
        empty_credits: dict = {"words_by": None, "music_by": None, "arranger": None}
        slides = self._make_slides(image_blob=b"\x89PNG\r\n")
        resolver.resolve(slides, empty_credits, canonical_title="some song")
        assert len(ocr_calls) == 0

    def test_ocr_not_called_when_use_ocr_false(self, monkeypatch):
        """OCR must not be called when use_ocr=False even if budget is available."""
        from worship_catalog.extractor import CreditResolver, OcrBudget
        import worship_catalog.extractor as ext_module

        ocr_calls: list[bytes] = []

        def fake_try_ocr(slides):
            ocr_calls.append(b"called")
            return "Words: Somebody"

        monkeypatch.setattr(ext_module, "_try_ocr_credits", fake_try_ocr)

        budget = OcrBudget(max_calls=10)
        resolver = CreditResolver(library_index={}, ocr_budget=budget, use_ocr=False)
        empty_credits: dict = {"words_by": None, "music_by": None, "arranger": None}
        slides = self._make_slides(image_blob=b"\x89PNG\r\n")
        resolver.resolve(slides, empty_credits, canonical_title="some song")
        assert len(ocr_calls) == 0


# ---------------------------------------------------------------------------
# Issue #57 — OCR budget refund must be logged at DEBUG level.
# ---------------------------------------------------------------------------


class TestOcrBudgetRefundLogging:
    """When OCR returns no credits, the budget refund must be logged — issue #57."""

    def _make_slides(self, image_blob: bytes | None = None) -> list[Slide]:
        images = [SlideImage(shape_id=1, blob=image_blob)] if image_blob else []
        return [
            Slide(
                index=0,
                hidden=False,
                text=SlideText(text_lines=["1 – Some Song", "PaperlessHymnal.com"]),
                images=images,
            )
        ]

    def test_refund_is_logged_at_debug_when_ocr_returns_nothing(
        self, monkeypatch, caplog
    ):
        """When _try_ocr_credits returns None, a DEBUG log mentioning 'refund' must appear."""
        from worship_catalog.extractor import CreditResolver, OcrBudget
        import worship_catalog.extractor as ext_module

        monkeypatch.setattr(ext_module, "_try_ocr_credits", lambda _slides: None)

        budget = OcrBudget(max_calls=5)
        resolver = CreditResolver(library_index={}, ocr_budget=budget, use_ocr=True)
        empty_credits: dict = {"words_by": None, "music_by": None, "arranger": None}
        slides = self._make_slides(image_blob=b"\x89PNG\r\n")

        with caplog.at_level(logging.DEBUG, logger="worship_catalog.extractor"):
            resolver.resolve(slides, empty_credits, canonical_title="some song")

        assert any(
            "refund" in r.message.lower() for r in caplog.records
        ), f"Expected a DEBUG 'refund' log, got: {[r.message for r in caplog.records]}"

    def test_refund_decrements_budget(self, monkeypatch):
        """After a refund the budget calls_made counter goes back to its prior value."""
        from worship_catalog.extractor import CreditResolver, OcrBudget
        import worship_catalog.extractor as ext_module

        monkeypatch.setattr(ext_module, "_try_ocr_credits", lambda _slides: None)

        budget = OcrBudget(max_calls=5)
        assert budget.calls_made == 0

        resolver = CreditResolver(library_index={}, ocr_budget=budget, use_ocr=True)
        empty_credits: dict = {"words_by": None, "music_by": None, "arranger": None}
        slides = self._make_slides(image_blob=b"\x89PNG\r\n")
        resolver.resolve(slides, empty_credits, canonical_title="some song")

        # After a refund, calls_made should be 0 again
        assert budget.calls_made == 0, (
            f"Expected calls_made=0 after refund, got {budget.calls_made}"
        )


# ---------------------------------------------------------------------------
# Issue #131 — OCR refund path: budget must not be consumed when credits are
# already present (skip path).
# ---------------------------------------------------------------------------


class TestOcrBudgetRefundPathIssue131:
    """Budget must not be consumed when OCR is skipped due to credits already present — issue #131."""

    def _make_slides(self, image_blob: bytes | None = None):
        from worship_catalog.pptx_reader import Slide, SlideImage, SlideText
        images = [SlideImage(shape_id=1, blob=image_blob)] if image_blob else []
        return [
            Slide(
                index=0,
                hidden=False,
                text=SlideText(text_lines=["1 – Great Is Thy Faithfulness"]),
                images=images,
            )
        ]

    def test_budget_not_consumed_when_credits_already_present(self, monkeypatch):
        """Budget slot must NOT be consumed when Step 1 returns due to existing credits."""
        from worship_catalog.extractor import CreditResolver, OcrBudget
        import worship_catalog.extractor as ext_module

        ocr_calls: list = []
        monkeypatch.setattr(ext_module, "_try_ocr_credits", lambda _: ocr_calls.append(1) or "Words: X")

        budget = OcrBudget(max_calls=10)
        assert budget.calls_made == 0

        resolver = CreditResolver(library_index=None, ocr_budget=budget, use_ocr=True)
        # Credits already present — Step 1 should return immediately
        credits_with_data = {"words_by": "John Newton", "music_by": None, "arranger": None}
        slides = self._make_slides(image_blob=b"\x89PNG\r\n")
        resolver.resolve(slides, credits_with_data, canonical_title="amazing grace")

        assert budget.calls_made == 0, (
            f"Budget must not be consumed when credits are already present, "
            f"got calls_made={budget.calls_made}"
        )
        assert len(ocr_calls) == 0, "OCR must not be called when credits already present"

    def test_budget_consumed_when_ocr_called(self, monkeypatch):
        """Budget IS consumed when OCR is actually called (no pre-existing credits)."""
        from worship_catalog.extractor import CreditResolver, OcrBudget
        import worship_catalog.extractor as ext_module

        monkeypatch.setattr(ext_module, "_try_ocr_credits", lambda _: "Words: Someone")

        budget = OcrBudget(max_calls=10)
        resolver = CreditResolver(library_index=None, ocr_budget=budget, use_ocr=True)
        empty_credits = {"words_by": None, "music_by": None, "arranger": None}
        slides = self._make_slides(image_blob=b"\x89PNG\r\n")
        resolver.resolve(slides, empty_credits, canonical_title="some song")

        assert budget.calls_made == 1, (
            f"Budget must be consumed when OCR is called, got calls_made={budget.calls_made}"
        )

    def test_remaining_accurate_after_skip(self, monkeypatch):
        """remaining count is unchanged after a skip (credits already present)."""
        from worship_catalog.extractor import CreditResolver, OcrBudget
        import worship_catalog.extractor as ext_module

        monkeypatch.setattr(ext_module, "_try_ocr_credits", lambda _: None)

        budget = OcrBudget(max_calls=5)
        resolver = CreditResolver(library_index=None, ocr_budget=budget, use_ocr=True)
        credits_with_data = {"words_by": "Somebody", "music_by": None, "arranger": None}
        slides = self._make_slides(image_blob=b"\x89PNG\r\n")
        resolver.resolve(slides, credits_with_data, canonical_title="amazing grace")

        assert budget.remaining == 5, (
            f"remaining must stay at 5 when OCR is skipped, got {budget.remaining}"
        )

    def test_budget_refunded_when_ocr_returns_nothing(self, monkeypatch):
        """Budget IS refunded when OCR is called but returns no credits."""
        from worship_catalog.extractor import CreditResolver, OcrBudget
        import worship_catalog.extractor as ext_module

        monkeypatch.setattr(ext_module, "_try_ocr_credits", lambda _: None)

        budget = OcrBudget(max_calls=5)
        resolver = CreditResolver(library_index=None, ocr_budget=budget, use_ocr=True)
        empty_credits = {"words_by": None, "music_by": None, "arranger": None}
        slides = self._make_slides(image_blob=b"\x89PNG\r\n")
        resolver.resolve(slides, empty_credits, canonical_title="some song")

        # OCR was called (consume), then returned nothing (refund)
        assert budget.calls_made == 0, (
            f"Budget must be refunded when OCR returns nothing, got calls_made={budget.calls_made}"
        )
        assert budget.remaining == 5


# ---------------------------------------------------------------------------
# Issue #146 — OcrBudget boundary conditions
# ---------------------------------------------------------------------------


class TestOcrBudgetBoundaryConditions:
    """OcrBudget boundary condition tests — issue #146."""

    def test_consume_returns_true_when_budget_remains(self):
        """consume() returns True when there is remaining budget."""
        from worship_catalog.extractor import OcrBudget
        budget = OcrBudget(max_calls=5)
        assert budget.consume() is True

    def test_consume_returns_false_when_at_cap(self):
        """consume() returns False when all budget is exhausted."""
        from worship_catalog.extractor import OcrBudget
        budget = OcrBudget(max_calls=2)
        budget.consume()
        budget.consume()
        # Now at cap
        assert budget.consume() is False

    def test_is_capped_true_when_exhausted(self):
        """is_capped is True after all budget is used."""
        from worship_catalog.extractor import OcrBudget
        budget = OcrBudget(max_calls=1)
        budget.consume()
        assert budget.is_capped is True

    def test_remaining_decrements_on_each_consume(self):
        """remaining decrements correctly on each consume() call."""
        from worship_catalog.extractor import OcrBudget
        budget = OcrBudget(max_calls=5)
        assert budget.remaining == 5
        budget.consume()
        assert budget.remaining == 4
        budget.consume()
        assert budget.remaining == 3

    def test_refund_increments_remaining(self):
        """refund() increments remaining back after a consume."""
        from worship_catalog.extractor import OcrBudget
        budget = OcrBudget(max_calls=5)
        budget.consume()
        assert budget.remaining == 4
        budget.refund()
        assert budget.remaining == 5

    def test_refund_does_not_exceed_max(self):
        """refund() when no calls have been made should not set remaining above max_calls."""
        from worship_catalog.extractor import OcrBudget
        budget = OcrBudget(max_calls=3)
        # calls_made is already 0 — refund should be a no-op
        budget.refund()
        assert budget.remaining == 3
        assert budget.calls_made == 0

    def test_max_calls_zero_immediately_capped(self):
        """OcrBudget(max_calls=0) is immediately capped and consume() returns False."""
        from worship_catalog.extractor import OcrBudget
        budget = OcrBudget(max_calls=0)
        assert budget.is_capped is True
        assert budget.consume() is False
        assert budget.remaining == 0

    def test_max_calls_none_unlimited_never_caps(self):
        """OcrBudget(max_calls=None) is never capped, consume() always returns True."""
        from worship_catalog.extractor import OcrBudget
        budget = OcrBudget(max_calls=None)
        assert budget.is_capped is False
        assert budget.remaining is None
        for _ in range(1000):
            assert budget.consume() is True
        assert budget.is_capped is False

    def test_remaining_is_none_for_unlimited_budget(self):
        """remaining is None for unlimited budget (max_calls=None)."""
        from worship_catalog.extractor import OcrBudget
        budget = OcrBudget(max_calls=None)
        assert budget.remaining is None


# ---------------------------------------------------------------------------
# Issue #128 — CreditResolver must use library_index when provided
# ---------------------------------------------------------------------------


class TestCreditResolverLibraryLookup:
    """CreditResolver must consult the library index when provided — issue #128."""

    def _make_slides(self) -> list:
        return [make_slide(index=0, lines=["Amazing Grace", "PaperlessHymnal.com"])]

    def test_library_index_is_used_when_provided(self, monkeypatch):
        """When library_index is provided and contains the song, credits come from it."""
        from worship_catalog.extractor import CreditResolver

        library = {"amazing grace": {"words_by": "John Newton", "music_by": None, "arranger": None}}
        resolver = CreditResolver(library_index=library, ocr_budget=None, use_ocr=False)
        empty_credits: dict = {"words_by": None, "music_by": None, "arranger": None}
        result = resolver.resolve(self._make_slides(), empty_credits, canonical_title="amazing grace")
        assert result["words_by"] == "John Newton"

    def test_library_index_is_skipped_when_none(self, monkeypatch):
        """When library_index=None, the library step is skipped entirely."""
        from worship_catalog.extractor import CreditResolver

        # Construct with None — should not find any credits
        resolver = CreditResolver(library_index=None, ocr_budget=None, use_ocr=False)
        empty_credits: dict = {"words_by": None, "music_by": None, "arranger": None}
        result = resolver.resolve(self._make_slides(), empty_credits, canonical_title="amazing grace")
        # Without library or OCR, credits remain empty
        assert result["words_by"] is None

    def test_credits_resolved_from_library_for_unknown_song(self):
        """A song with no text credits gets credits from the library index."""
        from worship_catalog.extractor import _create_song_occurrence

        slides = [make_slide(0, ["Amazing Grace", "PaperlessHymnal.com"])]
        library_index = {
            "amazing grace": {
                "words_by": "John Newton",
                "music_by": "Traditional",
                "arranger": None,
            }
        }
        result = _create_song_occurrence(
            1, "amazing grace", slides, library_index=library_index
        )
        assert result.words_by == "John Newton"
        assert result.music_by == "Traditional"

    def test_create_song_occurrence_without_library_leaves_credits_empty(self):
        """Without library_index, a song with no text credits has empty credits."""
        from worship_catalog.extractor import _create_song_occurrence

        slides = [make_slide(0, ["Amazing Grace", "PaperlessHymnal.com"])]
        result = _create_song_occurrence(1, "amazing grace", slides, library_index=None)
        assert result.words_by is None
        assert result.music_by is None


# ---------------------------------------------------------------------------
# Issue #129 — run_import() shared pipeline function
# ---------------------------------------------------------------------------


class TestRunImport:
    """run_import() provides a shared import pipeline callable by CLI and web — issue #129."""

    def _minimal_pptx_bytes(self) -> bytes:
        """Reuse the minimal PPTX factory from TestExtractSongsFileHash."""
        return TestExtractSongsFileHash()._minimal_pptx()

    def test_returns_import_result_with_song_count(self, tmp_path):
        """run_import returns an object with a songs_imported count."""
        from worship_catalog.import_service import ImportResult, run_import
        from worship_catalog.db import Database

        db_path = tmp_path / "test.db"
        db = Database(db_path)
        db.connect()
        db.init_schema()

        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(self._minimal_pptx_bytes())

        result = run_import(db, pptx_path)
        assert isinstance(result, ImportResult)
        assert isinstance(result.songs_imported, int)
        db.close()

    def test_works_with_no_library_index(self, tmp_path):
        """run_import works when library_index is not provided (defaults to None)."""
        from worship_catalog.import_service import run_import
        from worship_catalog.db import Database

        db_path = tmp_path / "test.db"
        db = Database(db_path)
        db.connect()
        db.init_schema()

        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(self._minimal_pptx_bytes())

        result = run_import(db, pptx_path, library_index=None)
        assert result is not None
        db.close()

    def test_works_with_no_ocr_budget(self, tmp_path):
        """run_import works when ocr_budget is not provided (defaults to None)."""
        from worship_catalog.import_service import run_import
        from worship_catalog.db import Database

        db_path = tmp_path / "test.db"
        db = Database(db_path)
        db.connect()
        db.init_schema()

        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(self._minimal_pptx_bytes())

        result = run_import(db, pptx_path, ocr_budget=None)
        assert result is not None
        db.close()

    def test_raises_on_invalid_pptx_path(self, tmp_path):
        """run_import raises an exception for a non-existent path."""
        from worship_catalog.import_service import run_import
        from worship_catalog.db import Database

        db_path = tmp_path / "test.db"
        db = Database(db_path)
        db.connect()
        db.init_schema()

        with pytest.raises(Exception):
            run_import(db, tmp_path / "nonexistent.pptx")
        db.close()
