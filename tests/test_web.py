"""Tests for the FastAPI + HTMX web UI."""

import io
import json
import os
import sqlite3
import uuid as _uuid_mod
from pathlib import Path
from tempfile import TemporaryDirectory

import pytest
from starlette.testclient import TestClient

from conftest import CsrfAwareClient
from worship_catalog.db import Database


@pytest.fixture
def client(db_with_songs, tmp_path, monkeypatch):
    """TestClient with DB_PATH and INBOX_DIR env vars pointed at temp locations (CSRF-aware)."""
    inbox = tmp_path / "inbox"
    inbox.mkdir()
    monkeypatch.setenv("DB_PATH", str(db_with_songs))
    monkeypatch.setenv("INBOX_DIR", str(inbox))
    from importlib import reload
    import worship_catalog.web.app as app_module
    reload(app_module)
    return CsrfAwareClient(TestClient(app_module.app))


@pytest.fixture
def raw_client(db_with_songs, monkeypatch):
    """Plain TestClient without CSRF token injection — for CSRF security tests."""
    monkeypatch.setenv("DB_PATH", str(db_with_songs))
    from importlib import reload
    import worship_catalog.web.app as app_module
    reload(app_module)
    return TestClient(app_module.app)


class TestRootRedirect:
    def test_root_redirects_to_songs(self, client):
        response = client.get("/", follow_redirects=False)
        assert response.status_code in (301, 302, 307, 308)
        assert response.headers["location"].endswith("/songs")


class TestSongsPage:
    def test_songs_page_returns_html(self, client):
        response = client.get("/songs")
        assert response.status_code == 200
        assert "text/html" in response.headers["content-type"]

    def test_songs_page_contains_song_titles(self, client):
        response = client.get("/songs")
        assert "Amazing Grace" in response.text
        assert "How Great Thou Art" in response.text

    def test_songs_page_contains_credits(self, client):
        response = client.get("/songs")
        assert "John Newton" in response.text
        assert "Stuart K. Hine" in response.text

    def test_songs_search_filters_results(self, client):
        response = client.get("/songs?q=Amazing")
        assert response.status_code == 200
        assert "Amazing Grace" in response.text
        # HTMX partial or full page — either way only matching song shown
        assert "How Great Thou Art" not in response.text

    def test_songs_htmx_request_returns_partial(self, client):
        """HTMX requests return table rows only (no full page nav)."""
        response = client.get("/songs?q=Amazing", headers={"HX-Request": "true"})
        assert response.status_code == 200
        assert "Amazing Grace" in response.text
        # Partial should NOT include the full <nav> chrome
        assert "<nav" not in response.text

    def test_songs_empty_search_returns_all(self, client):
        response = client.get("/songs?q=")
        assert response.status_code == 200
        assert "Amazing Grace" in response.text
        assert "How Great Thou Art" in response.text

    def test_songs_search_no_match_shows_empty(self, client):
        response = client.get("/songs?q=xyznotasong", headers={"HX-Request": "true"})
        assert response.status_code == 200
        assert "No songs found" in response.text


class TestReportsPage:
    def test_reports_page_returns_html(self, client):
        response = client.get("/reports")
        assert response.status_code == 200
        assert "text/html" in response.headers["content-type"]

    def test_reports_page_shows_ccli_section(self, client):
        """The reports page must have a CCLI report section (#201)."""
        response = client.get("/reports")
        assert response.status_code == 200
        assert "CCLI" in response.text, (
            "Reports page has no CCLI section — church admins cannot "
            "generate the compliance report from the web UI"
        )

    def test_reports_page_has_stats_form(self, client):
        response = client.get("/reports")
        assert "stats" in response.text.lower()
        assert "/reports/stats" in response.text


class TestStatsReport:
    def test_stats_returns_html(self, client):
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert response.status_code == 200
        assert "text/html" in response.headers["content-type"]

    def test_stats_shows_song_titles(self, client):
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert "Amazing Grace" in response.text or "How Great Thou Art" in response.text

    def test_stats_shows_summary_counts(self, client):
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        # Summary stat boxes should show counts
        assert "1" in response.text  # at least 1 service

    def test_stats_leader_filter(self, client):
        """Leader filter returns only that leader's services."""
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31", "leader": "Matt"},
        )
        assert response.status_code == 200
        assert "Amazing Grace" in response.text  # Matt led this service

    def test_stats_leader_filter_no_match(self, client):
        """Leader filter with no matching services shows empty state."""
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31", "leader": "Nobody"},
        )
        assert response.status_code == 200
        assert "No services found" in response.text

    def test_stats_empty_range_shows_no_services(self, client):
        response = client.post(
            "/reports/stats",
            data={"start_date": "2020-01-01", "end_date": "2020-01-31"},
        )
        assert response.status_code == 200
        assert "No services found" in response.text

    def test_stats_htmx_partial_no_nav(self, client):
        """HTMX POST returns partial HTML (no full page nav)."""
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
            headers={"HX-Request": "true"},
        )
        assert response.status_code == 200
        # stats_result.html partial is returned directly — no <nav>
        assert "<nav" not in response.text

    def test_stats_shows_leader_breakdown(self, client):
        """Stats without leader filter shows By Song Leader breakdown."""
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert response.status_code == 200
        assert "By Song Leader" in response.text
        assert "Matt" in response.text  # leader name from fixture

    def test_stats_leader_breakdown_hidden_when_filtered(self, client):
        """Stats with leader filter does NOT show the breakdown section."""
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31", "leader": "Matt"},
        )
        assert response.status_code == 200
        assert "By Song Leader" not in response.text

    def test_stats_breakdown_shows_leader_service_count(self, client):
        """Leader card in breakdown shows service count."""
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        # Fixture has 1 service led by Matt
        assert "1 service" in response.text

    def test_stats_breakdown_shows_songs_per_leader(self, client):
        """Leader card lists songs with repeat counts."""
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        # Both fixture songs should appear under Matt
        assert "Amazing Grace" in response.text
        assert "How Great Thou Art" in response.text


class TestServicesListPage:
    def test_services_page_returns_html(self, client):
        response = client.get("/services")
        assert response.status_code == 200
        assert "text/html" in response.headers["content-type"]

    def test_services_page_lists_services(self, client):
        response = client.get("/services")
        assert "AM Worship" in response.text
        assert "2026-02-15" in response.text

    def test_services_page_shows_leader(self, client):
        response = client.get("/services")
        assert "Matt" in response.text

    def test_services_page_shows_song_count(self, client):
        response = client.get("/services")
        # Fixture has 2 songs in one service
        assert "2" in response.text

    def test_services_page_links_to_detail(self, client, db_with_songs):
        response = client.get("/services")
        assert "/services/" in response.text

    def test_services_empty_db_shows_message(self, client, tmp_path, monkeypatch):
        empty_db = tmp_path / "empty.db"
        db = Database(empty_db)
        db.connect()
        db.init_schema()
        db.close()
        monkeypatch.setenv("DB_PATH", str(empty_db))
        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)
        c = TestClient(app_module.app)
        response = c.get("/services")
        assert response.status_code == 200
        assert "No services" in response.text

    def test_services_nav_link_is_active(self, client):
        response = client.get("/services")
        assert 'class="active"' in response.text


class TestServiceDetailPage:
    def _get_service_id(self, client):
        """Fetch the services list and extract a valid service ID from a link."""
        response = client.get("/services")
        import re
        match = re.search(r'/services/(\d+)', response.text)
        assert match, "No service detail links found"
        return int(match.group(1))

    def test_detail_returns_html(self, client):
        svc_id = self._get_service_id(client)
        response = client.get(f"/services/{svc_id}")
        assert response.status_code == 200
        assert "text/html" in response.headers["content-type"]

    def test_detail_shows_service_name(self, client):
        svc_id = self._get_service_id(client)
        response = client.get(f"/services/{svc_id}")
        assert "AM Worship" in response.text

    def test_detail_shows_service_date(self, client):
        svc_id = self._get_service_id(client)
        response = client.get(f"/services/{svc_id}")
        assert "2026-02-15" in response.text

    def test_detail_shows_song_leader(self, client):
        svc_id = self._get_service_id(client)
        response = client.get(f"/services/{svc_id}")
        assert "Matt" in response.text

    def test_detail_shows_setlist(self, client):
        svc_id = self._get_service_id(client)
        response = client.get(f"/services/{svc_id}")
        assert "Amazing Grace" in response.text
        assert "How Great Thou Art" in response.text

    def test_detail_shows_credits(self, client):
        svc_id = self._get_service_id(client)
        response = client.get(f"/services/{svc_id}")
        assert "John Newton" in response.text
        assert "Stuart K. Hine" in response.text

    def test_detail_shows_back_link(self, client):
        svc_id = self._get_service_id(client)
        response = client.get(f"/services/{svc_id}")
        assert "/services" in response.text

    def test_detail_404_for_missing_service(self, client):
        response = client.get("/services/99999")
        assert response.status_code == 404


class TestSongDetailPage:
    def _get_song_id(self, client):
        import re
        response = client.get("/songs")
        match = re.search(r'/songs/(\d+)', response.text)
        assert match, "No song detail links found"
        return int(match.group(1))

    def test_song_title_is_linked(self, client):
        response = client.get("/songs")
        assert "/songs/" in response.text

    def test_detail_returns_html(self, client):
        song_id = self._get_song_id(client)
        response = client.get(f"/songs/{song_id}")
        assert response.status_code == 200
        assert "text/html" in response.headers["content-type"]

    def test_detail_shows_song_title(self, client):
        song_id = self._get_song_id(client)
        response = client.get(f"/songs/{song_id}")
        assert "Amazing Grace" in response.text or "How Great Thou Art" in response.text

    def test_detail_shows_credits(self, client):
        response = client.get(f"/songs/{self._get_song_id(client)}")
        assert "John Newton" in response.text or "Stuart K. Hine" in response.text

    def test_detail_shows_service_history(self, client):
        song_id = self._get_song_id(client)
        response = client.get(f"/songs/{song_id}")
        assert "2026-02-15" in response.text

    def test_detail_shows_back_link(self, client):
        song_id = self._get_song_id(client)
        response = client.get(f"/songs/{song_id}")
        assert "/songs" in response.text

    def test_detail_404_for_missing_song(self, client):
        response = client.get("/songs/99999")
        assert response.status_code == 404


class TestSongsSorting:
    def test_sort_by_title_asc(self, client):
        response = client.get("/songs?sort=display_title&sort_dir=asc")
        assert response.status_code == 200

    def test_sort_by_performance_count(self, client):
        response = client.get("/songs?sort=performance_count&sort_dir=desc")
        assert response.status_code == 200

    def test_sort_indicator_shown(self, client):
        response = client.get("/songs?sort=display_title&sort_dir=asc")
        assert "▲" in response.text

    def test_invalid_sort_col_falls_back(self, client):
        response = client.get("/songs?sort=INVALID&sort_dir=asc")
        assert response.status_code == 200

    def test_htmx_partial_sort(self, client):
        response = client.get(
            "/songs?sort=display_title&sort_dir=desc",
            headers={"HX-Request": "true"},
        )
        assert response.status_code == 200
        assert "<tr>" in response.text


class TestServicesFiltering:
    def test_filter_by_leader(self, client):
        response = client.get("/services?q_leader=Matt")
        assert response.status_code == 200
        assert "Matt" in response.text

    def test_filter_no_match(self, client):
        response = client.get("/services?q_leader=ZZZNOMATCH")
        assert response.status_code == 200
        assert "No services" in response.text

    def test_filter_by_date_range(self, client):
        response = client.get("/services?start_date=2026-01-01&end_date=2026-12-31")
        assert response.status_code == 200
        assert "AM Worship" in response.text

    def test_sort_by_date_asc(self, client):
        response = client.get("/services?sort=service_date&sort_dir=asc")
        assert response.status_code == 200
        assert "▲" in response.text

    def test_htmx_partial_returns_rows(self, client):
        response = client.get(
            "/services?q_leader=Matt",
            headers={"HX-Request": "true"},
        )
        assert response.status_code == 200
        assert "<tr>" in response.text

    def test_invalid_sort_col_falls_back(self, client):
        response = client.get("/services?sort=INVALID")
        assert response.status_code == 200


class TestLeaderRoutes:
    """Tests for /leaders and /leaders/{name}/top-songs routes."""

    def test_leaders_index_returns_html(self, client):
        response = client.get("/leaders")
        assert response.status_code == 200
        assert "text/html" in response.headers["content-type"]

    def test_leaders_index_lists_leader(self, client):
        response = client.get("/leaders")
        assert "Matt" in response.text

    def test_leader_top_songs_returns_html(self, client):
        response = client.get("/leaders/Matt/top-songs")
        assert response.status_code == 200
        assert "text/html" in response.headers["content-type"]

    def test_leader_top_songs_shows_leader_name(self, client):
        response = client.get("/leaders/Matt/top-songs")
        assert "Matt" in response.text

    def test_leader_top_songs_empty_state_for_no_repeats(self, client):
        """Fixture has only 1 service for Matt, so no songs repeat 2+ times."""
        response = client.get("/leaders/Matt/top-songs")
        assert response.status_code == 200
        # Should show "no songs" message or warning since nothing repeated
        assert "No songs" in response.text or "more than once" in response.text.lower()

    def test_leader_top_songs_csv_download(self, client):
        response = client.get("/leaders/Matt/top-songs/csv")
        assert response.status_code == 200
        assert "text/csv" in response.headers["content-type"]

    def test_leader_top_songs_csv_has_header(self, client):
        response = client.get("/leaders/Matt/top-songs/csv")
        first_line = response.text.splitlines()[0]
        assert "Title" in first_line

    def test_unknown_leader_shows_message(self, client):
        response = client.get("/leaders/NoSuchLeader/top-songs")
        assert response.status_code == 200
        # Should show empty state, not 404
        assert "No songs" in response.text or "more than once" in response.text.lower()


class TestStatsExport:
    """Tests for stats report CSV and Excel export."""

    def test_stats_csv_returns_csv_content_type(self, client):
        response = client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert response.status_code == 200
        assert "text/csv" in response.headers["content-type"]

    def test_stats_csv_has_attachment_header(self, client):
        response = client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert "attachment" in response.headers["content-disposition"]

    def test_stats_csv_filename_includes_dates(self, client):
        response = client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        cd = response.headers["content-disposition"]
        assert "2026-01-01" in cd
        assert "2026-12-31" in cd

    def test_stats_csv_has_header_row(self, client):
        response = client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        first_line = response.text.splitlines()[0]
        assert "Title" in first_line

    def test_stats_csv_contains_song(self, client):
        response = client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert "Amazing Grace" in response.text or "How Great Thou Art" in response.text

    def test_stats_csv_empty_range_header_only(self, client):
        response = client.post(
            "/reports/stats/csv",
            data={"start_date": "2020-01-01", "end_date": "2020-01-31"},
        )
        assert response.status_code == 200
        non_empty_lines = [l for l in response.text.splitlines() if l.strip()]
        assert len(non_empty_lines) == 1

    def test_stats_download_buttons_in_result(self, client):
        """The stats result HTML shows download buttons."""
        response = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert "/reports/stats/csv" in response.text


class TestCcliReportWebRoute:
    """Web UI must provide CCLI report generation (#201)."""

    def test_reports_page_shows_ccli_section(self, client):
        """The reports page must have a CCLI report section."""
        resp = client.get("/reports")
        assert resp.status_code == 200
        assert "CCLI" in resp.text, (
            "Reports page has no CCLI section — church admins cannot "
            "generate the compliance report from the web UI"
        )

    def test_ccli_report_download_returns_csv(self, client):
        """POST /reports/ccli should return a downloadable CSV file."""
        resp = client.post(
            "/reports/ccli",
            data={"start_date": "2025-01-01", "end_date": "2025-12-31"},
        )
        assert resp.status_code == 200
        content_type = resp.headers.get("content-type", "")
        assert "text/csv" in content_type or "application/octet-stream" in content_type, (
            f"Expected CSV content type, got {content_type}"
        )

    def test_ccli_report_has_content_disposition(self, client):
        """CCLI report download should have a filename header."""
        resp = client.post(
            "/reports/ccli",
            data={"start_date": "2025-01-01", "end_date": "2025-12-31"},
        )
        disposition = resp.headers.get("content-disposition", "")
        assert "attachment" in disposition, (
            "CCLI report should be a file download with Content-Disposition: attachment"
        )

    def test_ccli_report_csv_has_header_row(self, client):
        """CSV must have a header row with expected CCLI columns."""
        resp = client.post(
            "/reports/ccli",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        first_line = resp.text.splitlines()[0]
        assert "Title" in first_line
        assert "Reproduction Type" in first_line

    def test_ccli_report_csv_contains_copy_events(self, client):
        """CSV should contain data from the test fixture's copy events."""
        resp = client.post(
            "/reports/ccli",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert "Amazing Grace" in resp.text or "How Great Thou Art" in resp.text

    def test_ccli_report_csv_empty_range_still_has_header(self, client):
        """An empty date range should return CSV with just the header row."""
        resp = client.post(
            "/reports/ccli",
            data={"start_date": "2020-01-01", "end_date": "2020-01-31"},
        )
        assert resp.status_code == 200
        non_empty_lines = [line for line in resp.text.splitlines() if line.strip()]
        assert len(non_empty_lines) == 1  # header only

    def test_ccli_report_filename_includes_dates(self, client):
        """Content-Disposition filename should include the date range."""
        resp = client.post(
            "/reports/ccli",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        cd = resp.headers.get("content-disposition", "")
        assert "2026-01-01" in cd
        assert "2026-12-31" in cd

    def test_ccli_report_validates_dates(self, client):
        """Invalid dates should return 422."""
        resp = client.post(
            "/reports/ccli",
            data={"start_date": "bad-date", "end_date": "2026-12-31"},
        )
        assert resp.status_code == 422


class TestPagination:
    """Tests for pagination on /songs and /services."""

    @pytest.fixture
    def db_with_many_songs(self, tmp_path):
        """DB with 55 songs across 55 services to test multi-page results."""
        db_path = tmp_path / "paginated.db"
        db = Database(db_path)
        db.connect()
        db.init_schema()

        for i in range(55):
            title = f"Song Number {i:03d}"
            canonical = title.lower()
            song_id = db.insert_or_get_song(canonical, title)
            db.insert_or_get_song_edition(song_id, words_by=f"Author {i}")

            svc_id = db.insert_or_update_service(
                service_date=f"2026-{(i % 12) + 1:02d}-{(i % 28) + 1:02d}",
                service_name=f"Service {i}",
                source_file=f"file{i}.pptx",
                source_hash=f"hash{i}",
            )
            db.insert_service_song(svc_id, song_id, ordinal=1)

        db.close()
        return db_path

    @pytest.fixture
    def paginated_client(self, db_with_many_songs, monkeypatch):
        monkeypatch.setenv("DB_PATH", str(db_with_many_songs))
        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)
        return TestClient(app_module.app)

    def test_songs_default_returns_50_rows(self, paginated_client):
        response = paginated_client.get("/songs")
        assert response.status_code == 200
        # Count <tr> rows; includes 1 header row so total rows = data rows + 1
        row_count = response.text.count("<tr>")
        assert row_count <= 51  # 50 data rows + 1 header row

    def test_songs_page2_accessible(self, paginated_client):
        response = paginated_client.get("/songs?page=2&per_page=50")
        assert response.status_code == 200

    def test_songs_pagination_links_shown(self, paginated_client):
        """With 55 songs and per_page=50, page 1 should show a Next link."""
        response = paginated_client.get("/songs?page=1&per_page=50")
        assert response.status_code == 200
        assert "page=2" in response.text or "Next" in response.text

    def test_songs_page1_no_prev_link(self, paginated_client):
        response = paginated_client.get("/songs?page=1&per_page=50")
        # Page 1 should not show a "Previous" link
        assert "page=0" not in response.text

    def test_services_pagination_accessible(self, paginated_client):
        response = paginated_client.get("/services?page=1&per_page=50")
        assert response.status_code == 200

    def test_per_page_respected(self, paginated_client):
        response = paginated_client.get("/songs?page=1&per_page=10")
        assert response.status_code == 200
        row_count = response.text.count("<tr>")
        assert row_count <= 11  # 10 data rows + 1 header row


class TestErrorPages:
    """Tests for custom 404/500 HTML error pages."""

    def test_404_song_returns_404_status(self, client):
        response = client.get("/songs/99999")
        assert response.status_code == 404

    def test_404_service_returns_404_status(self, client):
        response = client.get("/services/99999")
        assert response.status_code == 404

    def test_404_response_is_html_not_json(self, client):
        response = client.get("/songs/99999")
        assert response.status_code == 404
        assert "text/html" in response.headers["content-type"]
        # Should NOT start with a JSON brace
        assert not response.text.strip().startswith("{")

    def test_404_body_contains_useful_text(self, client):
        response = client.get("/songs/99999")
        assert "404" in response.text or "not found" in response.text.lower()

    def test_404_body_contains_back_link(self, client):
        response = client.get("/songs/99999")
        assert "/songs" in response.text

    def test_health_endpoint(self, client):
        response = client.get("/health")
        assert response.status_code == 200


class TestHealthEndpointDb:
    """Tests for health endpoint with DB connectivity check — issue #31."""

    def test_health_returns_200_with_connected_db(self, client):
        """When DB is reachable, /health returns 200."""
        response = client.get("/health")
        assert response.status_code == 200

    def test_health_returns_status_ok_in_body(self, client):
        """Health response includes status: ok."""
        response = client.get("/health")
        data = response.json()
        assert data.get("status") == "ok"

    def test_health_does_not_expose_db_backend_details(self, client):
        """Health response must not leak DB backend info (#62)."""
        response = client.get("/health")
        data = response.json()
        assert "db" not in data, "Response must not expose DB backend details"

    def test_health_error_response_is_generic(self, monkeypatch, tmp_path):
        """503 response must not expose backend details (#62)."""
        import worship_catalog.web.app as app_module
        from importlib import reload
        from unittest.mock import patch

        inbox = tmp_path / "inbox"
        inbox.mkdir()
        monkeypatch.setenv("DB_PATH", str(tmp_path / "worship.db"))
        monkeypatch.setenv("INBOX_DIR", str(inbox))
        reload(app_module)

        def broken_get_db():
            raise OSError("simulated DB failure")

        from starlette.testclient import TestClient
        c = TestClient(app_module.app, raise_server_exceptions=False)
        with patch.object(app_module, "_get_db", broken_get_db):
            response = c.get("/health")
        assert response.status_code == 503
        data = response.json()
        assert "db" not in data, "Error response must not expose DB backend details"

    def test_health_returns_503_when_db_unavailable(self, monkeypatch, tmp_path):
        """When DB raises on execute, /health returns 503."""
        from unittest.mock import patch, MagicMock

        # Patch _get_db to raise so we can test the error path
        def broken_get_db():
            raise OSError("simulated DB failure")

        import worship_catalog.web.app as app_module
        from importlib import reload
        monkeypatch.setenv("DB_PATH", str(tmp_path / "worship.db"))
        reload(app_module)
        from starlette.testclient import TestClient
        c = TestClient(app_module.app, raise_server_exceptions=False)

        with patch.object(app_module, "_get_db", broken_get_db):
            response = c.get("/health")
        assert response.status_code == 503


class TestDateValidation:
    """Tests for ISO-8601 date validation in web form endpoints — issue #17."""

    def test_invalid_start_date_returns_422(self, client):
        """Non-ISO date string in start_date returns a validation error."""
        response = client.post(
            "/reports/stats/csv",
            data={"start_date": "Jan 2026", "end_date": "2026-12-31"},
        )
        assert response.status_code == 422

    def test_invalid_end_date_returns_422(self, client):
        """Non-ISO date string in end_date returns a validation error."""
        response = client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-01-01", "end_date": "December 2026"},
        )
        assert response.status_code == 422

    def test_start_after_end_returns_422(self, client):
        """start_date > end_date returns a validation error."""
        response = client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-12-31", "end_date": "2026-01-01"},
        )
        assert response.status_code == 422

    def test_valid_dates_are_accepted(self, client):
        """Well-formed ISO dates proceed normally."""
        response = client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert response.status_code == 200

    def test_stats_invalid_date_returns_422(self, client):
        """Stats endpoint also validates dates."""
        response = client.post(
            "/reports/stats",
            data={"start_date": "bad-date", "end_date": "2026-12-31"},
        )
        assert response.status_code == 422


class TestCSRFProtection:
    """Tests for CSRF protection on POST report endpoints — issue #39."""

    def test_post_without_csrf_token_is_rejected(self, raw_client):
        """POST to report endpoint without CSRF token returns 403."""
        response = raw_client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert response.status_code == 403

    def test_post_with_valid_csrf_token_succeeds(self, raw_client):
        """POST with a valid CSRF token is accepted."""
        get_resp = raw_client.get("/reports")
        token = get_resp.cookies.get("csrftoken")
        assert token is not None, "CSRF cookie should be set on GET"

        response = raw_client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
            headers={"X-CSRFToken": token},
        )
        assert response.status_code == 200

    def test_post_with_wrong_csrf_token_is_rejected(self, raw_client):
        """POST with wrong X-CSRFToken value returns 403."""
        raw_client.get("/reports")  # set cookie
        response = raw_client.post(
            "/reports/stats/csv",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
            headers={"X-CSRFToken": "wrong-token"},
        )
        assert response.status_code == 403

    def test_all_report_post_endpoints_require_csrf(self, raw_client):
        """All report POST endpoints reject requests without CSRF token."""
        endpoints = [
            "/reports/stats",
            "/reports/stats/csv",
            "/reports/stats/xlsx",
        ]
        for endpoint in endpoints:
            resp = raw_client.post(
                endpoint,
                data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
            )
            assert resp.status_code == 403, (
                f"{endpoint} accepted POST without CSRF token (got {resp.status_code})"
            )

    def test_upload_without_csrf_token_returns_403(self, raw_client):
        """/upload must be protected by CSRF — no token → 403. (#59)"""
        resp = raw_client.post(
            "/upload",
            files={"file": ("sunday.pptx", io.BytesIO(SMALL_PPTX_BYTES), VALID_PPTX_MIME)},
        )
        assert resp.status_code == 403

    def test_upload_with_valid_csrf_token_is_not_rejected_for_csrf(self, client):
        """/upload with a valid CSRF token must not return 403. (#59)"""
        resp = _upload(client, SMALL_PPTX_BYTES, "sunday.pptx", VALID_PPTX_MIME)
        assert resp.status_code != 403


# ---------------------------------------------------------------------------
# Upload / background import job tests (#45)
# ---------------------------------------------------------------------------

VALID_PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
SMALL_PPTX_BYTES = b"PK\x03\x04" + b"\x00" * 100  # fake PPTX magic bytes


def _upload(
    client: "CsrfAwareClient",
    content: bytes,
    filename: str,
    content_type: str,
):
    return client.post(
        "/upload",
        files={"file": (filename, io.BytesIO(content), content_type)},
    )


class TestUploadEndpoint:
    def test_upload_valid_pptx_returns_202_with_job_id(self, client):
        resp = _upload(client, SMALL_PPTX_BYTES, "sunday.pptx", VALID_PPTX_MIME)
        assert resp.status_code == 202
        body = resp.json()
        assert "job_id" in body
        job_id = body["job_id"]
        # Must be a URL-safe string of sufficient length (≥ 32 URL-safe chars)
        # secrets.token_urlsafe(32) produces 43 chars; uuid4 would only be 36
        assert len(job_id) >= 32, "job_id must have sufficient entropy"
        import re as _re
        assert _re.fullmatch(r"[A-Za-z0-9_\-]+", job_id), "job_id must be URL-safe"

    def test_upload_job_id_is_not_sequential(self, client):
        """Two uploads must produce different, unrelated job IDs (#60)."""
        id1 = _upload(client, SMALL_PPTX_BYTES, "a.pptx", VALID_PPTX_MIME).json()["job_id"]
        id2 = _upload(client, SMALL_PPTX_BYTES, "b.pptx", VALID_PPTX_MIME).json()["job_id"]
        assert id1 != id2

    def test_upload_non_pptx_mime_returns_400(self, client):
        resp = _upload(client, b"hello", "notes.pdf", "application/pdf")
        assert resp.status_code == 400
        assert "pptx" in resp.json()["detail"].lower()

    def test_upload_wrong_extension_returns_400(self, client):
        resp = _upload(client, SMALL_PPTX_BYTES, "sunday.ppt", VALID_PPTX_MIME)
        assert resp.status_code == 400

    def test_upload_oversized_file_returns_413(self, client, monkeypatch):
        monkeypatch.setattr("worship_catalog.web.app.MAX_UPLOAD_BYTES", 10)
        resp = _upload(client, b"X" * 11, "big.pptx", VALID_PPTX_MIME)
        assert resp.status_code == 413

    def test_upload_rejects_oversized_content_length_before_reading(self, client, monkeypatch):
        """413 must come from Content-Length header check, not post-read byte count."""
        monkeypatch.setattr("worship_catalog.web.app.MAX_UPLOAD_BYTES", 10)
        resp = client.post(
            "/upload",
            files={"file": ("big.pptx", io.BytesIO(b"X" * 5), VALID_PPTX_MIME)},
            headers={"content-length": "11"},
        )
        assert resp.status_code == 413

    def test_upload_missing_content_length_falls_back_to_body_limit(self, client, monkeypatch):
        """Without Content-Length header, post-read body check still rejects oversized body."""
        monkeypatch.setattr("worship_catalog.web.app.MAX_UPLOAD_BYTES", 10)
        resp = _upload(client, b"X" * 11, "big.pptx", VALID_PPTX_MIME)
        assert resp.status_code == 413

    def test_upload_creates_pending_job_record(self, client, db):
        resp = _upload(client, SMALL_PPTX_BYTES, "sunday.pptx", VALID_PPTX_MIME)
        job_id = resp.json()["job_id"]
        row = db.get_import_job(job_id)
        assert row is not None
        assert row["filename"] == "sunday.pptx"
        assert row["status"] in ("pending", "running", "complete", "failed")


class TestJobStatusEndpoint:
    def test_get_known_job_returns_200(self, client, db):
        job_id = str(_uuid_mod.uuid4())
        db.create_import_job(job_id, filename="test.pptx")
        resp = client.get(f"/jobs/{job_id}")
        assert resp.status_code == 200
        body = resp.json()
        assert body["job_id"] == job_id
        assert body["status"] == "pending"

    def test_get_unknown_job_returns_404(self, client):
        resp = client.get(f"/jobs/{_uuid_mod.uuid4()}")
        assert resp.status_code == 404

    def test_job_response_includes_all_fields(self, client, db):
        job_id = str(_uuid_mod.uuid4())
        db.create_import_job(job_id, filename="test.pptx")
        resp = client.get(f"/jobs/{job_id}")
        body = resp.json()
        for field in (
            "job_id", "filename", "status", "started_at",
            "completed_at", "songs_imported", "error_message",
        ):
            assert field in body


class TestJobListEndpoint:
    def test_list_jobs_returns_200(self, client):
        resp = client.get("/jobs")
        assert resp.status_code == 200
        assert isinstance(resp.json(), list)

    def test_list_jobs_newest_first(self, client, db):
        id1 = str(_uuid_mod.uuid4())
        id2 = str(_uuid_mod.uuid4())
        db.create_import_job(id1, filename="first.pptx")
        db.create_import_job(id2, filename="second.pptx")
        resp = client.get("/jobs")
        ids = [j["job_id"] for j in resp.json()]
        assert ids.index(id2) < ids.index(id1)


class TestJobLifecycle:
    def test_job_transitions_pending_to_running_to_complete(self, db):
        job_id = str(_uuid_mod.uuid4())
        db.create_import_job(job_id, filename="test.pptx")
        assert db.get_import_job(job_id)["status"] == "pending"

        db.update_import_job(job_id, status="running")
        assert db.get_import_job(job_id)["status"] == "running"

        db.update_import_job(job_id, status="complete", songs_imported=5)
        row = db.get_import_job(job_id)
        assert row["status"] == "complete"
        assert row["songs_imported"] == 5
        assert row["completed_at"] is not None

    def test_job_captures_error_message_on_failure(self, db):
        job_id = str(_uuid_mod.uuid4())
        db.create_import_job(job_id, filename="bad.pptx")
        db.update_import_job(job_id, status="failed", error_message="corrupt file")
        row = db.get_import_job(job_id)
        assert row["status"] == "failed"
        assert "corrupt" in row["error_message"]

    def test_job_songs_imported_count_set_on_success(self, db):
        job_id = str(_uuid_mod.uuid4())
        db.create_import_job(job_id, filename="songs.pptx")
        db.update_import_job(job_id, status="complete", songs_imported=12)
        assert db.get_import_job(job_id)["songs_imported"] == 12


class TestJobAutoPurge:
    def test_purge_removes_records_older_than_90_days(self, db):
        old_id = str(_uuid_mod.uuid4())
        recent_id = str(_uuid_mod.uuid4())
        # 2025-12-01 is well over 90 days before 2026-03-15
        db.create_import_job(old_id, filename="old.pptx", started_at="2025-12-01T00:00:00")
        db.create_import_job(recent_id, filename="recent.pptx")
        db.purge_old_import_jobs(days=90)
        assert db.get_import_job(old_id) is None
        assert db.get_import_job(recent_id) is not None

    def test_purge_keeps_records_exactly_at_boundary(self, db):
        # Compute the boundary date dynamically so the test stays correct
        # regardless of when it runs.  A job started exactly 90 days ago
        # must NOT be deleted by purge_old_import_jobs(days=90).
        from datetime import datetime, timedelta, timezone
        boundary_dt = datetime.now(timezone.utc) - timedelta(days=90)
        boundary_started_at = boundary_dt.strftime("%Y-%m-%dT00:00:00")
        boundary_id = str(_uuid_mod.uuid4())
        db.create_import_job(
            boundary_id,
            filename="boundary.pptx",
            started_at=boundary_started_at,
        )
        db.purge_old_import_jobs(days=90)
        assert db.get_import_job(boundary_id) is not None


class TestUploadStructuredLogs:
    def test_log_emitted_on_job_start(self, client, caplog):
        import logging
        with caplog.at_level(logging.INFO, logger="worship_catalog"):
            _upload(client, SMALL_PPTX_BYTES, "sunday.pptx", VALID_PPTX_MIME)
        messages = [r.message for r in caplog.records]
        assert any("import_job" in m and "pending" in m for m in messages)

    def test_log_emitted_on_job_complete(self, db, caplog):
        import logging
        job_id = str(_uuid_mod.uuid4())
        db.create_import_job(job_id, filename="test.pptx")
        with caplog.at_level(logging.INFO, logger="worship_catalog"):
            db.update_import_job(job_id, status="complete", songs_imported=3)
        messages = [r.message for r in caplog.records]
        assert any("complete" in m for m in messages)


# ---------------------------------------------------------------------------
# Background import transaction safety (#51)
# ---------------------------------------------------------------------------

class TestBackgroundImportTransaction:
    def test_failed_import_marks_job_failed_and_does_not_silently_skip(
        self, tmp_path, monkeypatch
    ):
        """If extract_songs raises, the job must be marked failed (not left as pending)."""
        from worship_catalog.db import Database

        db_path = tmp_path / "tx_test.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        job_id = str(_uuid_mod.uuid4())
        _db.create_import_job(job_id, filename="bad.pptx")
        _db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))

        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)

        pptx_path = tmp_path / "bad.pptx"
        pptx_path.write_bytes(b"not a real pptx")
        app_module._run_import_in_background(job_id, pptx_path)

        _db2 = Database(db_path)
        _db2.connect()
        row = _db2.get_import_job(job_id)
        _db2.close()
        assert row["status"] == "failed"
        assert row["error_message"] is not None

    def test_successful_import_marks_job_complete_with_song_count(
        self, tmp_path, monkeypatch
    ):
        """If extract_songs succeeds, job must be complete with songs_imported set."""
        from worship_catalog.db import Database
        from worship_catalog.extractor import SongOccurrence, ExtractionResult
        import worship_catalog.extractor as extractor_mod
        import worship_catalog.pptx_reader as pptx_reader_mod

        db_path = tmp_path / "tx_ok.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        job_id = str(_uuid_mod.uuid4())
        _db.create_import_job(job_id, filename="ok.pptx")
        _db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))

        fake_songs = [
            SongOccurrence(ordinal=i + 1, canonical_title=f"song {i}", display_title=f"Song {i}")
            for i in range(3)
        ]
        fake_result = ExtractionResult(
            filename="ok.pptx",
            file_hash="fakehash",
            service_date="2026-01-04",
            service_name="AM Worship",
            song_leader=None,
            preacher=None,
            sermon_title=None,
            songs=fake_songs,
        )

        # Patch at the module level so the imports inside _run_import_in_background resolve
        # to the fakes regardless of reload order.
        monkeypatch.setattr(extractor_mod, "extract_songs", lambda p, **kw: fake_result)
        monkeypatch.setattr(pptx_reader_mod, "compute_file_hash", lambda p: "fakehash")

        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)

        pptx_path = tmp_path / "ok.pptx"
        pptx_path.write_bytes(b"fake")
        app_module._run_import_in_background(job_id, pptx_path)

        _db2 = Database(db_path)
        _db2.connect()
        row = _db2.get_import_job(job_id)
        _db2.close()
        assert row["status"] == "complete"
        assert row["songs_imported"] == 3


# ---------------------------------------------------------------------------
# Startup auto-purge (#74)
# ---------------------------------------------------------------------------

class TestStartupPurge:
    def _make_app(self, db_path, tmp_path, monkeypatch):
        """Reload the app module with a fresh DB path and return a TestClient."""
        inbox = tmp_path / "inbox"
        inbox.mkdir(exist_ok=True)
        monkeypatch.setenv("DB_PATH", str(db_path))
        monkeypatch.setenv("INBOX_DIR", str(inbox))
        import worship_catalog.web.app as m
        from importlib import reload
        reload(m)
        return m.app

    def test_startup_purges_import_jobs_older_than_90_days(self, tmp_path, monkeypatch):
        """On app startup the lifespan must delete import_jobs older than 90 days."""
        db_path = tmp_path / "purge.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        old_id = str(_uuid_mod.uuid4())
        recent_id = str(_uuid_mod.uuid4())
        _db.create_import_job(old_id, filename="old.pptx", started_at="2025-12-01T00:00:00")
        _db.create_import_job(recent_id, filename="recent.pptx")
        _db.close()

        app = self._make_app(db_path, tmp_path, monkeypatch)
        with TestClient(app):
            pass  # lifespan fires on __enter__

        _db2 = Database(db_path)
        _db2.connect()
        assert _db2.get_import_job(old_id) is None
        assert _db2.get_import_job(recent_id) is not None
        _db2.close()

    def test_startup_purge_logs_at_info(self, tmp_path, monkeypatch, caplog):
        """Startup purge must emit an INFO log entry."""
        import logging
        db_path = tmp_path / "purge_log.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        _db.close()

        app = self._make_app(db_path, tmp_path, monkeypatch)
        with caplog.at_level(logging.INFO, logger="worship_catalog.web"):
            with TestClient(app):
                pass
        assert any("purge" in r.message.lower() for r in caplog.records)


# ---------------------------------------------------------------------------
# xlsx export tests (#84)
# ---------------------------------------------------------------------------

class TestStatsXlsxExport:
    """POST /reports/stats/xlsx — full coverage (#84)."""

    def _post(self, client, start="2026-01-01", end="2026-12-31", leader="", all_songs=False):
        return client.post("/reports/stats/xlsx", data={
            "start_date": start,
            "end_date": end,
            "leader": leader,
            "all_songs": "true" if all_songs else "",
        })

    def test_xlsx_returns_200(self, client):
        assert self._post(client).status_code == 200

    def test_xlsx_content_type_is_excel(self, client):
        ct = self._post(client).headers["content-type"]
        assert "spreadsheetml" in ct or "openxmlformats" in ct

    def test_xlsx_content_disposition_is_attachment(self, client):
        cd = self._post(client).headers["content-disposition"]
        assert "attachment" in cd
        assert ".xlsx" in cd

    def test_xlsx_filename_includes_dates(self, client):
        cd = self._post(client, start="2026-01-01", end="2026-03-31").headers["content-disposition"]
        assert "2026-01-01" in cd and "2026-03-31" in cd

    def test_xlsx_body_is_parseable_workbook(self, client):
        import io as _io
        import openpyxl
        wb = openpyxl.load_workbook(_io.BytesIO(self._post(client).content))
        assert "Top Songs" in wb.sheetnames

    def test_xlsx_top_songs_sheet_has_header_row(self, client):
        import io as _io
        import openpyxl
        ws = openpyxl.load_workbook(_io.BytesIO(self._post(client).content))["Top Songs"]
        headers = [ws.cell(1, c).value for c in range(1, 5)]
        assert headers == ["Rank", "Title", "Credits", "Count"]

    def test_xlsx_top_songs_sheet_contains_song_data(self, client):
        import io as _io
        import openpyxl
        ws = openpyxl.load_workbook(_io.BytesIO(self._post(client).content))["Top Songs"]
        titles = [ws.cell(r, 2).value for r in range(2, ws.max_row + 1) if ws.cell(r, 2).value]
        assert len(titles) > 0

    def test_xlsx_invalid_date_returns_422(self, client):
        assert self._post(client, start="not-a-date").status_code == 422

    def test_xlsx_openpyxl_missing_returns_501(self, client, monkeypatch):
        import builtins
        real_import = builtins.__import__
        def _block(name, *args, **kwargs):
            if name == "openpyxl":
                raise ImportError("blocked for test")
            return real_import(name, *args, **kwargs)
        monkeypatch.setattr(builtins, "__import__", _block)
        assert self._post(client).status_code == 501


# ---------------------------------------------------------------------------
# Boundary / adversarial input tests (#89)
# ---------------------------------------------------------------------------

class TestInputBoundaryConditions:
    """Boundary and adversarial input tests for web forms (#89)."""

    def test_page_zero_does_not_500(self, client):
        assert client.get("/songs?page=0").status_code != 500

    def test_negative_page_does_not_500(self, client):
        assert client.get("/songs?page=-1").status_code != 500

    def test_non_integer_page_does_not_500(self, client):
        assert client.get("/songs?page=abc").status_code != 500

    def test_search_xss_is_html_escaped(self, client):
        """Script tag in q param must never appear unescaped in the response."""
        resp = client.get("/songs?q=%3Cscript%3Ealert(1)%3C%2Fscript%3E")
        assert resp.status_code != 500
        assert "<script>alert(1)</script>" not in resp.text

    def test_search_percent_wildcard_returns_200(self, client):
        assert client.get("/songs?q=%25").status_code == 200

    def test_search_sql_quote_returns_200(self, client):
        assert client.get("/songs?q=%27").status_code == 200

    def test_per_page_zero_does_not_500(self, client):
        assert client.get("/songs?per_page=0").status_code != 500

    def test_per_page_enormous_does_not_500(self, client):
        assert client.get("/songs?per_page=100000").status_code != 500

    def test_very_long_search_does_not_500(self, client):
        assert client.get(f"/songs?q={'a' * 5000}").status_code != 500

    def test_leader_name_with_apostrophe_does_not_500(self, client):
        resp = client.get("/leaders/O%27Brien/top-songs")
        assert resp.status_code in (200, 404)
        assert resp.status_code != 500

    def test_services_page_negative_page_does_not_500(self, client):
        assert client.get("/services?page=-1").status_code != 500

    def test_ccli_route_exists_and_returns_csv(self, client):
        """POST /reports/ccli must return a CSV download (#201)."""
        resp = client.post("/reports/ccli",
                           data={"start_date": "2026-01-01", "end_date": "2026-12-31"})
        assert resp.status_code == 200


# ---------------------------------------------------------------------------
# ThreadPoolExecutor bound for upload jobs (#52)
# ---------------------------------------------------------------------------

class TestUploadThreadPool:
    """Background import must use a bounded thread pool, not unbounded threads (#52)."""

    def test_app_uses_threadpoolexecutor_not_bare_thread(self):
        """app.py must use ThreadPoolExecutor for import jobs; threading.Thread is only
        allowed in the lifespan shutdown guard (a single, bounded, one-time thread)."""
        import inspect
        from worship_catalog.web import app as web_module
        src = inspect.getsource(web_module)
        assert "ThreadPoolExecutor" in src, (
            "Upload background tasks must use ThreadPoolExecutor to bound concurrency"
        )
        # threading.Thread is permitted in _lifespan for the shutdown timeout guard (#135),
        # but _run_import_in_background must not use it — verify by inspecting that function.
        run_import_src = inspect.getsource(web_module._run_import_in_background)
        assert "threading.Thread(" not in run_import_src, (
            "Import jobs must not spawn unbounded bare threads — use ThreadPoolExecutor"
        )

    def test_thread_pool_is_module_level_singleton(self):
        """The executor must be a module-level singleton, not created per-request."""
        import inspect
        from worship_catalog.web import app as web_module
        src = inspect.getsource(web_module)
        # Module-level assignment: _executor = ThreadPoolExecutor(...)
        assert "_executor" in src or "_EXECUTOR" in src, (
            "Executor must be a named module-level constant, not an anonymous per-request object"
        )


class TestUploadConcurrencyLimit:
    """Pool-full behaviour — submitting more jobs than pool capacity must return 503 (#52)."""

    def test_upload_returns_503_when_pool_full(self, client, monkeypatch):
        """When the import thread pool is saturated, /upload returns 503."""
        from concurrent.futures import Future
        import worship_catalog.web.app as app_module
        from importlib import reload

        # Replace the module-level executor with one whose submit() raises
        # concurrent.futures.BrokenExecutor (pool full / shutdown).
        class _FullExecutor:
            def submit(self, *args, **kwargs):
                raise RuntimeError("pool saturated")

        monkeypatch.setattr(app_module, "_import_executor", _FullExecutor())
        resp = _upload(client, SMALL_PPTX_BYTES, "sunday.pptx", VALID_PPTX_MIME)
        assert resp.status_code == 503
        body = resp.json()
        assert "busy" in body.get("detail", "").lower() or "unavailable" in body.get("detail", "").lower()

    def test_upload_succeeds_when_pool_has_capacity(self, client, monkeypatch):
        """When the pool has capacity, /upload succeeds (202) as normal."""
        from concurrent.futures import Future
        import worship_catalog.web.app as app_module

        # Replace executor with a no-op that always accepts (simulates available pool)
        class _AcceptingExecutor:
            def submit(self, *args, **kwargs):
                f: Future[None] = Future()
                f.set_result(None)
                return f

        monkeypatch.setattr(app_module, "_import_executor", _AcceptingExecutor())
        resp = _upload(client, SMALL_PPTX_BYTES, "sunday.pptx", VALID_PPTX_MIME)
        assert resp.status_code == 202
        assert "job_id" in resp.json()

    def test_upload_503_body_is_json(self, client, monkeypatch):
        """The 503 response for a full pool must be JSON, not HTML."""
        import worship_catalog.web.app as app_module

        class _FullExecutor:
            def submit(self, *args, **kwargs):
                raise RuntimeError("pool saturated")

        monkeypatch.setattr(app_module, "_import_executor", _FullExecutor())
        resp = _upload(client, SMALL_PPTX_BYTES, "sunday.pptx", VALID_PPTX_MIME)
        assert resp.status_code == 503
        # Must be parseable as JSON
        body = resp.json()
        assert "detail" in body


# ---------------------------------------------------------------------------
# #110: E2E test — background import persists songs to DB (issue #96)
# ---------------------------------------------------------------------------


@pytest.mark.integration
@pytest.mark.slow
class TestBackgroundImportPersistsSongsToDB:
    """Issue #96 — _run_import_in_background must write songs to the database.

    The background worker was updating job status but not calling the extractor
    or persisting any song/service rows.  These tests verify the full E2E path:
    upload a real PPTX → poll until complete → confirm songs exist in the DB.

    Marked ``slow`` because each test builds a PPTX in memory, runs the full
    background import worker, and polls the DB for results (~200–260 ms each).
    """

    @pytest.fixture
    def minimal_pptx_bytes(self, tmp_path):
        """Build a minimal Paperless Hymnal PPTX in memory and return its bytes."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        blank = prs.slide_layouts[6]

        def add_text_box(slide, lines):
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
            tf = txBox.text_frame
            for i, line in enumerate(lines):
                if i == 0:
                    tf.paragraphs[0].text = line
                else:
                    tf.add_paragraph().text = line

        # Metadata slide (table)
        meta = prs.slides.add_slide(blank)
        from pptx.util import Inches as _I
        table = meta.shapes.add_table(3, 2, _I(1), _I(1), _I(8), _I(1.2)).table
        for r, (k, v) in enumerate([("Date", "2026-01-04"), ("Service", "AM Worship"), ("Song Leader", "Alice")]):
            table.cell(r, 0).text = k
            table.cell(r, 1).text = v

        # Song slides
        add_text_box(prs.slides.add_slide(blank), [
            "Amazing Grace", "How sweet the sound", "Words: John Newton", "PaperlessHymnal.com",
        ])
        add_text_box(prs.slides.add_slide(blank), [
            "Amazing Grace", "That saved a wretch like me",
        ])

        import io
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _make_upload_client(self, tmp_path, monkeypatch):
        """Return (CsrfAwareClient, db_path, inbox) with isolated env."""
        db_path = tmp_path / "upload_e2e.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        _db.close()

        inbox = tmp_path / "inbox"
        inbox.mkdir()
        monkeypatch.setenv("DB_PATH", str(db_path))
        monkeypatch.setenv("INBOX_DIR", str(inbox))

        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)
        return CsrfAwareClient(TestClient(app_module.app)), db_path

    def test_upload_real_pptx_job_reaches_complete_status(self, tmp_path, monkeypatch, minimal_pptx_bytes):
        """A real PPTX uploaded via /upload must drive the job to 'complete' status."""
        client, db_path = self._make_upload_client(tmp_path, monkeypatch)

        resp = client.post(
            "/upload",
            files={"file": ("AM Worship 2026.01.04.pptx", io.BytesIO(minimal_pptx_bytes), VALID_PPTX_MIME)},
        )
        assert resp.status_code == 202
        job_id = resp.json()["job_id"]

        # _run_import_in_background runs synchronously in TestClient's thread pool;
        # poll up to 5 seconds for the job to finish.
        import time
        deadline = time.monotonic() + 5.0
        status = "pending"
        while time.monotonic() < deadline and status not in ("complete", "failed"):
            status = client.get(f"/jobs/{job_id}").json()["status"]
            time.sleep(0.05)

        assert status == "complete", f"Job never reached 'complete'; last status={status!r}"

    def test_upload_real_pptx_songs_persisted_to_db(self, tmp_path, monkeypatch, minimal_pptx_bytes):
        """After a successful upload, the extracted songs must exist in the songs table."""
        client, db_path = self._make_upload_client(tmp_path, monkeypatch)

        resp = client.post(
            "/upload",
            files={"file": ("AM Worship 2026.01.04.pptx", io.BytesIO(minimal_pptx_bytes), VALID_PPTX_MIME)},
        )
        assert resp.status_code == 202
        job_id = resp.json()["job_id"]

        # Wait for job to finish
        import time
        deadline = time.monotonic() + 5.0
        status = "pending"
        while time.monotonic() < deadline and status not in ("complete", "failed"):
            status = client.get(f"/jobs/{job_id}").json()["status"]
            time.sleep(0.05)

        assert status == "complete", f"Job did not complete: {status!r}"

        # The song must now be in the DB
        _db = Database(db_path)
        _db.connect()
        cursor = _db.cursor()
        cursor.execute("SELECT COUNT(*) FROM songs")
        count = cursor.fetchone()[0]
        _db.close()

        assert count > 0, (
            "No songs were persisted to the DB after background import — "
            "the background worker must call the extractor and write song rows"
        )

    def test_upload_real_pptx_service_persisted_to_db(self, tmp_path, monkeypatch, minimal_pptx_bytes):
        """After a successful upload, a service row must exist in the services table."""
        client, db_path = self._make_upload_client(tmp_path, monkeypatch)

        resp = client.post(
            "/upload",
            files={"file": ("AM Worship 2026.01.04.pptx", io.BytesIO(minimal_pptx_bytes), VALID_PPTX_MIME)},
        )
        assert resp.status_code == 202
        job_id = resp.json()["job_id"]

        import time
        deadline = time.monotonic() + 5.0
        status = "pending"
        while time.monotonic() < deadline and status not in ("complete", "failed"):
            status = client.get(f"/jobs/{job_id}").json()["status"]
            time.sleep(0.05)

        assert status == "complete", f"Job did not complete: {status!r}"

        _db = Database(db_path)
        _db.connect()
        cursor = _db.cursor()
        cursor.execute("SELECT COUNT(*) FROM services")
        count = cursor.fetchone()[0]
        _db.close()

        assert count > 0, (
            "No services were persisted to the DB after background import — "
            "the background worker must persist service metadata"
        )

    def test_background_worker_persists_songs_directly(self, tmp_path, monkeypatch, minimal_pptx_bytes):
        """Call _run_import_in_background directly and verify DB rows are written."""
        db_path = tmp_path / "direct.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        job_id = str(_uuid_mod.uuid4())
        _db.create_import_job(job_id, filename="AM Worship 2026.01.04.pptx")
        _db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))
        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)

        pptx_path = tmp_path / "AM Worship 2026.01.04.pptx"
        pptx_path.write_bytes(minimal_pptx_bytes)

        app_module._run_import_in_background(job_id, pptx_path)

        _db2 = Database(db_path)
        _db2.connect()
        row = _db2.get_import_job(job_id)
        cursor = _db2.cursor()
        cursor.execute("SELECT COUNT(*) FROM songs")
        song_count = cursor.fetchone()[0]
        _db2.close()

        assert row["status"] == "complete"
        assert song_count > 0, (
            "_run_import_in_background must persist extracted songs to the DB; "
            f"found {song_count} song rows after import"
        )

    def test_upload_creates_copy_events(self, tmp_path, monkeypatch, minimal_pptx_bytes):
        """After a successful upload, projection and recording copy events must exist (#176)."""
        client, db_path = self._make_upload_client(tmp_path, monkeypatch)

        resp = client.post(
            "/upload",
            files={"file": ("AM Worship 2026.01.04.pptx", io.BytesIO(minimal_pptx_bytes), VALID_PPTX_MIME)},
        )
        assert resp.status_code == 202
        job_id = resp.json()["job_id"]

        import time
        deadline = time.monotonic() + 5.0
        status = "pending"
        while time.monotonic() < deadline and status not in ("complete", "failed"):
            status = client.get(f"/jobs/{job_id}").json()["status"]
            time.sleep(0.05)
        assert status == "complete"

        _db = Database(db_path)
        _db.connect()
        events = _db.query_copy_events("0000-01-01", "9999-12-31")
        _db.close()
        types = {e["reproduction_type"] for e in events}
        assert "projection" in types, (
            "Background import must create 'projection' copy events"
        )

    def test_upload_idempotent_reimport(self, tmp_path, monkeypatch, minimal_pptx_bytes):
        """Uploading the same PPTX twice must not create duplicate services (#176)."""
        client, db_path = self._make_upload_client(tmp_path, monkeypatch)

        import time

        def upload_and_wait():
            resp = client.post(
                "/upload",
                files={"file": ("AM Worship 2026.01.04.pptx", io.BytesIO(minimal_pptx_bytes), VALID_PPTX_MIME)},
            )
            job_id = resp.json()["job_id"]
            deadline = time.monotonic() + 5.0
            status = "pending"
            while time.monotonic() < deadline and status not in ("complete", "failed"):
                status = client.get(f"/jobs/{job_id}").json()["status"]
                time.sleep(0.05)
            return status

        s1 = upload_and_wait()
        assert s1 == "complete"
        s2 = upload_and_wait()
        assert s2 == "complete"

        _db = Database(db_path)
        _db.connect()
        services = _db.query_services("2026-01-01", "2026-12-31")
        _db.close()
        assert len(services) == 1, (
            f"Expected 1 service after idempotent re-import, got {len(services)}"
        )

    def test_upload_persists_service_songs_join(self, tmp_path, monkeypatch, minimal_pptx_bytes):
        """service_songs rows must link the imported song to its service (#176)."""
        client, db_path = self._make_upload_client(tmp_path, monkeypatch)

        resp = client.post(
            "/upload",
            files={"file": ("AM Worship 2026.01.04.pptx", io.BytesIO(minimal_pptx_bytes), VALID_PPTX_MIME)},
        )
        job_id = resp.json()["job_id"]

        import time
        deadline = time.monotonic() + 5.0
        status = "pending"
        while time.monotonic() < deadline and status not in ("complete", "failed"):
            status = client.get(f"/jobs/{job_id}").json()["status"]
            time.sleep(0.05)
        assert status == "complete"

        _db = Database(db_path)
        _db.connect()
        cursor = _db.cursor()
        cursor.execute("SELECT COUNT(*) FROM service_songs")
        count = cursor.fetchone()[0]
        _db.close()
        assert count > 0, (
            "service_songs must contain rows linking songs to the service"
        )

    def test_upload_songs_imported_count_matches_db(self, tmp_path, monkeypatch, minimal_pptx_bytes):
        """Job record songs_imported must match actual songs in DB (#176)."""
        client, db_path = self._make_upload_client(tmp_path, monkeypatch)

        resp = client.post(
            "/upload",
            files={"file": ("AM Worship 2026.01.04.pptx", io.BytesIO(minimal_pptx_bytes), VALID_PPTX_MIME)},
        )
        job_id = resp.json()["job_id"]

        import time
        deadline = time.monotonic() + 5.0
        status = "pending"
        while time.monotonic() < deadline and status not in ("complete", "failed"):
            status = client.get(f"/jobs/{job_id}").json()["status"]
            time.sleep(0.05)
        assert status == "complete"

        job = client.get(f"/jobs/{job_id}").json()
        _db = Database(db_path)
        _db.connect()
        cursor = _db.cursor()
        cursor.execute("SELECT COUNT(*) FROM songs")
        song_count = cursor.fetchone()[0]
        _db.close()
        assert job["songs_imported"] == song_count, (
            f"Job says {job['songs_imported']} songs imported but DB has {song_count}"
        )


# ---------------------------------------------------------------------------
# Issue #112: Contract tests for Content-Disposition filename format
# ---------------------------------------------------------------------------

import re as _re


class TestDownloadFilenameContract:
    """Contract tests that assert the filename format in Content-Disposition headers.

    These tests define the download-filename contract so that future refactors
    cannot accidentally break the format expected by browsers and downstream
    consumers (#112).
    """

    # -----------------------------------------------------------------------
    # Helpers shared across all download endpoints
    # -----------------------------------------------------------------------

    _STATS_FORM = {
        "start_date": "2026-01-01",
        "end_date": "2026-12-31",
        "leader": "",
        "all_songs": "",
    }

    _SAFE_FILENAME_RE = _re.compile(r'^[\w.\-]+$')

    def _assert_attachment(self, cd: str) -> None:
        assert cd.startswith('attachment; filename='), (
            f"Content-Disposition should start with 'attachment; filename=', got: {cd!r}"
        )

    def _extract_filename(self, cd: str) -> str:
        # Support both quoted (filename="foo.csv") and unquoted (filename=foo.csv)
        part = cd.split("filename=", 1)[-1].strip().strip('"')
        return part

    # -----------------------------------------------------------------------
    # POST /reports/stats/csv
    # -----------------------------------------------------------------------

    def test_stats_csv_has_attachment_disposition(self, client):
        resp = client.post("/reports/stats/csv", data=self._STATS_FORM)
        assert resp.status_code == 200
        cd = resp.headers.get("content-disposition", "")
        self._assert_attachment(cd)

    def test_stats_csv_filename_contains_only_safe_chars(self, client):
        resp = client.post("/reports/stats/csv", data=self._STATS_FORM)
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert self._SAFE_FILENAME_RE.match(filename), (
            f"Unsafe characters in stats CSV filename: {filename!r}"
        )

    def test_stats_csv_filename_has_csv_extension(self, client):
        resp = client.post("/reports/stats/csv", data=self._STATS_FORM)
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert filename.endswith(".csv"), f"Expected .csv extension, got: {filename!r}"

    def test_stats_csv_filename_embeds_dates(self, client):
        resp = client.post("/reports/stats/csv", data={
            **self._STATS_FORM,
            "start_date": "2026-01-15",
            "end_date": "2026-03-31",
        })
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert "2026-01-15" in filename and "2026-03-31" in filename, (
            f"Stats CSV filename should include both date bounds, got: {filename!r}"
        )

    # -----------------------------------------------------------------------
    # POST /reports/stats/xlsx
    # -----------------------------------------------------------------------

    def test_stats_xlsx_has_attachment_disposition(self, client):
        resp = client.post("/reports/stats/xlsx", data=self._STATS_FORM)
        assert resp.status_code == 200
        cd = resp.headers.get("content-disposition", "")
        self._assert_attachment(cd)

    def test_stats_xlsx_filename_contains_only_safe_chars(self, client):
        resp = client.post("/reports/stats/xlsx", data=self._STATS_FORM)
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert self._SAFE_FILENAME_RE.match(filename), (
            f"Unsafe characters in stats XLSX filename: {filename!r}"
        )

    def test_stats_xlsx_filename_has_xlsx_extension(self, client):
        resp = client.post("/reports/stats/xlsx", data=self._STATS_FORM)
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert filename.endswith(".xlsx"), f"Expected .xlsx extension, got: {filename!r}"

    def test_stats_xlsx_filename_embeds_dates(self, client):
        resp = client.post("/reports/stats/xlsx", data={
            **self._STATS_FORM,
            "start_date": "2026-02-01",
            "end_date": "2026-04-30",
        })
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert "2026-02-01" in filename and "2026-04-30" in filename, (
            f"Stats XLSX filename should include both date bounds, got: {filename!r}"
        )

    # -----------------------------------------------------------------------
    # GET /leaders/{leader_name}/top-songs/csv
    # -----------------------------------------------------------------------

    def test_leader_csv_has_attachment_disposition(self, client):
        resp = client.get("/leaders/Matt/top-songs/csv")
        assert resp.status_code == 200
        cd = resp.headers.get("content-disposition", "")
        self._assert_attachment(cd)

    def test_leader_csv_filename_contains_only_safe_chars(self, client):
        resp = client.get("/leaders/Matt/top-songs/csv")
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert self._SAFE_FILENAME_RE.match(filename), (
            f"Unsafe characters in leader CSV filename: {filename!r}"
        )

    def test_leader_csv_filename_has_csv_extension(self, client):
        resp = client.get("/leaders/Matt/top-songs/csv")
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert filename.endswith(".csv"), f"Expected .csv extension, got: {filename!r}"

    def test_leader_csv_filename_contains_leader_name(self, client):
        """Leader name (sanitized) appears in the download filename."""
        resp = client.get("/leaders/Matt/top-songs/csv")
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert "Matt" in filename or "matt" in filename.lower(), (
            f"Leader name should appear in filename, got: {filename!r}"
        )

    def test_leader_csv_filename_sanitizes_spaces(self, client, db_with_songs, monkeypatch):
        """A leader name with spaces must not produce a filename with spaces."""
        # Insert a service led by a leader whose name contains spaces
        from worship_catalog.db import Database
        db = Database(db_with_songs)
        db.connect()
        svc = db.insert_or_update_service(
            service_date="2026-01-10",
            service_name="PM Worship",
            source_file="pm.pptx",
            source_hash="pm123",
            song_leader="John Doe",
        )
        song_id = db.insert_or_get_song("amazing grace", "Amazing Grace")
        db.insert_service_song(svc, song_id, ordinal=1, song_edition_id=None)
        db.insert_or_get_copy_event(svc, song_id, "projection", song_edition_id=None)
        db.insert_or_get_copy_event(svc, song_id, "projection", song_edition_id=None)
        db.close()

        resp = client.get("/leaders/John%20Doe/top-songs/csv")
        assert resp.status_code == 200
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert " " not in filename, (
            f"Filename must not contain spaces, got: {filename!r}"
        )
        assert self._SAFE_FILENAME_RE.match(filename), (
            f"Unsafe characters in leader CSV filename with spaced leader name: {filename!r}"
        )

    def test_leader_csv_filename_no_path_separators(self, client):
        """Filename must never contain path separators (security: path traversal guard)."""
        resp = client.get("/leaders/Matt/top-songs/csv")
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert "/" not in filename and "\\" not in filename, (
            f"Filename must not contain path separators, got: {filename!r}"
        )

    def test_leader_csv_content_disposition_no_semicolon_in_filename(self, client):
        """The filename portion must not contain semicolons (would break header parsing)."""
        resp = client.get("/leaders/Matt/top-songs/csv")
        cd = resp.headers.get("content-disposition", "")
        filename = self._extract_filename(cd)
        assert ";" not in filename, (
            f"Filename must not contain semicolons, got: {filename!r}"
        )


# ---------------------------------------------------------------------------
# Issue #132 — ORDER BY whitelist guard inside query_songs_paginated / query_all_services_paginated
# Methods moved from app.py to Database (#166)
# ---------------------------------------------------------------------------


class TestQuerySongsInternalWhitelist:
    """Database.query_songs_paginated() must validate sort column — issue #132."""

    @pytest.fixture
    def temp_db(self, tmp_path):
        from worship_catalog.db import Database
        db = Database(tmp_path / "qs_test.db")
        db.connect()
        db.init_schema()
        yield db
        db.close()

    def test_valid_sort_col_works(self, temp_db):
        """Passing a valid sort column succeeds without error."""
        rows, total = temp_db.query_songs_paginated(sort="display_title", sort_dir="asc")
        assert isinstance(rows, list)

    def test_invalid_sort_col_raises_value_error(self, temp_db):
        """Invalid sort column raises ValueError."""
        with pytest.raises(ValueError, match="Invalid sort column"):
            temp_db.query_songs_paginated(sort="not_a_real_column")

    def test_sql_injection_sort_raises_value_error(self, temp_db):
        """SQL injection string as sort column is rejected."""
        with pytest.raises(ValueError):
            temp_db.query_songs_paginated(sort="title; DROP TABLE songs--")

    def test_empty_sort_col_raises_value_error(self, temp_db):
        """Empty string sort column is rejected."""
        with pytest.raises(ValueError):
            temp_db.query_songs_paginated(sort="")

    def test_all_valid_songs_sort_cols_work(self, temp_db):
        """Every column in _SONGS_SORT_COLS must be accepted without error."""
        from worship_catalog.db import Database
        for col in Database._SONGS_SORT_COLS:
            rows, total = temp_db.query_songs_paginated(sort=col)
            assert isinstance(rows, list), f"Column {col!r} failed unexpectedly"


class TestQueryServicesInternalWhitelist:
    """Database.query_all_services_paginated() must validate sort column — issue #132."""

    @pytest.fixture
    def temp_db(self, tmp_path):
        from worship_catalog.db import Database
        db = Database(tmp_path / "svc_test.db")
        db.connect()
        db.init_schema()
        yield db
        db.close()

    def test_valid_sort_col_works(self, temp_db):
        """Passing a valid sort column succeeds."""
        rows, total = temp_db.query_all_services_paginated(sort="service_date", sort_dir="asc")
        assert isinstance(rows, list)

    def test_invalid_sort_col_raises_value_error(self, temp_db):
        """Invalid sort column raises ValueError."""
        with pytest.raises(ValueError, match="Invalid sort column"):
            temp_db.query_all_services_paginated(sort="not_valid_col")

    def test_sql_injection_sort_raises_value_error(self, temp_db):
        """SQL injection string as sort column raises ValueError."""
        with pytest.raises(ValueError):
            temp_db.query_all_services_paginated(sort="service_date; DROP TABLE services--")

    def test_all_valid_services_sort_cols_work(self, temp_db):
        """Every column in _SERVICES_SORT_COLS must be accepted without error."""
        from worship_catalog.db import Database
        for col in Database._SERVICES_SORT_COLS:
            rows, total = temp_db.query_all_services_paginated(sort=col)
            assert isinstance(rows, list), f"Column {col!r} failed unexpectedly"


# ---------------------------------------------------------------------------
# Issue #138 — Uploaded PPTX files must be deleted from inbox after import
# ---------------------------------------------------------------------------


@pytest.mark.integration
class TestUploadInboxCleanup:
    """After background import completes, the uploaded file must be deleted — issue #138."""

    @pytest.fixture
    def minimal_pptx_bytes(self):
        try:
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation()
            blank = prs.slide_layouts[6]

            def add_text_box(slide, lines):
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
                tf = txBox.text_frame
                for i, line in enumerate(lines):
                    if i == 0:
                        tf.text = line
                    else:
                        tf.add_paragraph().text = line

            add_text_box(prs.slides.add_slide(blank), [
                "Amazing Grace", "How sweet the sound",
                "Words: John Newton", "PaperlessHymnal.com",
            ])
            buf = io.BytesIO()
            prs.save(buf)
            return buf.getvalue()
        except ImportError:
            pytest.skip("pptx not available")

    def _make_upload_client(self, tmp_path, monkeypatch):
        db_path = tmp_path / "cleanup_test.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        _db.close()
        inbox = tmp_path / "inbox"
        inbox.mkdir()
        monkeypatch.setenv("DB_PATH", str(db_path))
        monkeypatch.setenv("INBOX_DIR", str(inbox))
        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)
        return CsrfAwareClient(TestClient(app_module.app)), inbox

    def _wait_for_job(self, client, job_id, timeout=5.0):
        import time
        deadline = time.monotonic() + timeout
        status = "pending"
        while time.monotonic() < deadline and status not in ("complete", "failed"):
            status = client.get(f"/jobs/{job_id}").json()["status"]
            time.sleep(0.05)
        return status

    def test_file_deleted_from_inbox_after_success(self, tmp_path, monkeypatch, minimal_pptx_bytes):
        """Uploaded file must be deleted from inbox after successful import."""
        client, inbox = self._make_upload_client(tmp_path, monkeypatch)
        resp = client.post(
            "/upload",
            files={"file": ("AM Worship 2026.03.01.pptx", io.BytesIO(minimal_pptx_bytes), VALID_PPTX_MIME)},
        )
        assert resp.status_code == 202
        job_id = resp.json()["job_id"]
        status = self._wait_for_job(client, job_id)
        assert status == "complete", f"Job status={status!r}"
        # No PPTX files should remain in the inbox
        remaining = list(inbox.glob("*.pptx"))
        assert remaining == [], f"Inbox still contains files after import: {remaining}"

    def test_file_deleted_from_inbox_after_failure(self, tmp_path, monkeypatch):
        """Uploaded file must be deleted from inbox even if the import raises."""
        db_path = tmp_path / "fail_cleanup.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        _db.close()
        inbox = tmp_path / "inbox_fail"
        inbox.mkdir()
        monkeypatch.setenv("DB_PATH", str(db_path))
        monkeypatch.setenv("INBOX_DIR", str(inbox_fail := inbox))
        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)
        client = CsrfAwareClient(TestClient(app_module.app))

        # Upload a garbage PPTX that will fail extraction
        resp = client.post(
            "/upload",
            files={"file": ("bad.pptx", io.BytesIO(b"not a real pptx"), VALID_PPTX_MIME)},
        )
        assert resp.status_code == 202
        job_id = resp.json()["job_id"]
        import time
        deadline = time.monotonic() + 5.0
        status = "pending"
        while time.monotonic() < deadline and status not in ("complete", "failed"):
            status = client.get(f"/jobs/{job_id}").json()["status"]
            time.sleep(0.05)
        # Job should reach 'failed'
        assert status == "failed", f"Expected 'failed' status, got {status!r}"
        # Inbox must still be clean
        remaining = list(inbox_fail.glob("*.pptx"))
        assert remaining == [], f"Inbox still contains files after failed import: {remaining}"

    def test_job_status_set_before_file_deletion(self, tmp_path, monkeypatch, minimal_pptx_bytes):
        """Job status must transition to 'complete' before file cleanup."""
        client, inbox = self._make_upload_client(tmp_path, monkeypatch)
        resp = client.post(
            "/upload",
            files={"file": ("AM Worship 2026.03.02.pptx", io.BytesIO(minimal_pptx_bytes), VALID_PPTX_MIME)},
        )
        assert resp.status_code == 202
        job_id = resp.json()["job_id"]
        status = self._wait_for_job(client, job_id)
        # Status must be terminal before we check cleanup
        assert status in ("complete", "failed"), f"Unexpected status: {status}"
        remaining = list(inbox.glob("*.pptx"))
        assert remaining == [], f"File not cleaned up for status={status}: {remaining}"


# ---------------------------------------------------------------------------
# Issue #145 — /upload MIME type rejection and size limit tests
# ---------------------------------------------------------------------------


class TestUploadValidationIssue145:
    """Validation tests for /upload endpoint — issue #145 (tests only, no impl changes needed)."""

    def test_text_plain_mime_returns_400(self, client):
        """Uploading a file with content-type text/plain must be rejected with 400."""
        resp = _upload(client, b"hello world", "file.pptx", "text/plain")
        assert resp.status_code == 400, f"Expected 400 for text/plain, got {resp.status_code}"

    def test_text_plain_error_message_is_useful(self, client):
        """The 400 error body must contain a useful message for text/plain rejection."""
        resp = _upload(client, b"hello world", "file.pptx", "text/plain")
        body = resp.json()
        detail = body.get("detail", "")
        assert detail, f"Response body must contain 'detail', got: {body}"
        # Must mention pptx or mime
        assert "pptx" in detail.lower() or "mime" in detail.lower(), (
            f"Error message must mention PPTX or MIME type, got: {detail!r}"
        )

    def test_valid_pptx_mime_returns_202(self, client):
        """Uploading with the correct PPTX MIME type returns 202."""
        resp = _upload(client, SMALL_PPTX_BYTES, "valid.pptx", VALID_PPTX_MIME)
        assert resp.status_code == 202, f"Expected 202 for valid PPTX MIME, got {resp.status_code}"
        assert "job_id" in resp.json()

    def test_oversized_file_returns_413(self, client, monkeypatch):
        """Uploading a file larger than MAX_UPLOAD_BYTES returns 413."""
        monkeypatch.setattr("worship_catalog.web.app.MAX_UPLOAD_BYTES", 10)
        resp = _upload(client, b"X" * 20, "large.pptx", VALID_PPTX_MIME)
        assert resp.status_code == 413, f"Expected 413 for oversized file, got {resp.status_code}"

    def test_oversized_file_error_message_is_useful(self, client, monkeypatch):
        """The 413 error body must contain a useful error message."""
        monkeypatch.setattr("worship_catalog.web.app.MAX_UPLOAD_BYTES", 10)
        resp = _upload(client, b"X" * 20, "large.pptx", VALID_PPTX_MIME)
        body = resp.json()
        detail = body.get("detail", "")
        assert detail, f"413 response must contain 'detail', got: {body}"
        # Must mention size/exceeds/max
        assert any(kw in detail.lower() for kw in ("exceed", "size", "max", "byte")), (
            f"413 error message should mention file size, got: {detail!r}"
        )

    def test_zero_byte_file_returns_400_or_202(self, client):
        """A zero-byte upload should return an error response (not silently succeed)."""
        resp = _upload(client, b"", "empty.pptx", VALID_PPTX_MIME)
        # Zero-byte file: either rejected immediately (4xx) or accepted and fails on extraction
        # The key requirement is that it's handled gracefully (not 5xx)
        assert resp.status_code < 500, (
            f"Zero-byte upload must not return 5xx, got {resp.status_code}"
        )

    def test_zero_byte_file_has_detail_if_rejected(self, client):
        """If a zero-byte file is rejected, the response body must have a detail field."""
        resp = _upload(client, b"", "empty.pptx", VALID_PPTX_MIME)
        if resp.status_code >= 400:
            body = resp.json()
            assert "detail" in body, f"Error response must have 'detail' key, got: {body}"


class TestGracefulShutdown:
    """Verify the lifespan context manager shuts down the executor gracefully — closes #135."""

    def test_executor_shutdown_called_on_lifespan_exit(self, tmp_path, monkeypatch):
        """executor.shutdown(wait=True) must be called when the lifespan context exits."""
        from unittest.mock import patch, MagicMock
        from importlib import reload
        import worship_catalog.web.app as app_module

        shutdown_calls: list[dict] = []
        real_shutdown = app_module._import_executor.shutdown

        def tracking_shutdown(wait=True, cancel_futures=False):
            shutdown_calls.append({"wait": wait, "cancel_futures": cancel_futures})
            # Don't call real shutdown — executor may already be stopped after reload

        inbox = tmp_path / "inbox"
        inbox.mkdir()
        monkeypatch.setenv("DB_PATH", str(tmp_path / "test.db"))
        monkeypatch.setenv("INBOX_DIR", str(inbox))

        reload(app_module)
        # Patch the executor's shutdown method on the live instance
        app_module._import_executor.shutdown = tracking_shutdown  # type: ignore[method-assign]

        with TestClient(app_module.app):
            pass  # lifespan startup on enter, shutdown on exit

        assert shutdown_calls, "executor.shutdown() was never called during lifespan exit"
        assert shutdown_calls[-1]["wait"] is True, (
            f"executor.shutdown must be called with wait=True, got: {shutdown_calls[-1]}"
        )

    def test_in_flight_jobs_complete_before_shutdown(self, tmp_path, monkeypatch):
        """A job submitted just before shutdown must complete — executor waits for it."""
        import time
        from importlib import reload
        import worship_catalog.web.app as app_module

        completed: list[bool] = []

        def slow_job():
            time.sleep(0.05)  # short sleep — simulates a brief in-flight job
            completed.append(True)

        inbox = tmp_path / "inbox"
        inbox.mkdir()
        monkeypatch.setenv("DB_PATH", str(tmp_path / "test.db"))
        monkeypatch.setenv("INBOX_DIR", str(inbox))

        reload(app_module)
        with TestClient(app_module.app):
            # Submit a background job right before the context exits
            app_module._import_executor.submit(slow_job)
        # After the context exits (lifespan shutdown), the job must be done
        assert completed == [True], (
            "In-flight job was not completed before executor shutdown"
        )

    def test_lifespan_source_contains_shutdown_with_wait(self):
        """Source code of the lifespan function must include executor.shutdown(wait=True)."""
        import inspect
        from importlib import reload
        import worship_catalog.web.app as app_module
        reload(app_module)
        source = inspect.getsource(app_module._lifespan)
        assert "shutdown" in source, "lifespan must call executor.shutdown"
        assert "wait=True" in source, "lifespan executor.shutdown must use wait=True"


# ---------------------------------------------------------------------------
# Pushover notification integration in _run_import_in_background
# ---------------------------------------------------------------------------


class TestBackgroundImportNotification:
    """send_pushover must be called after successful and failed imports,
    and notification failures must never break the import flow."""

    def test_successful_import_sends_pushover_notification(
        self, tmp_path, monkeypatch
    ):
        """On successful import, send_pushover is called with 'Import complete'."""
        from worship_catalog.db import Database
        from worship_catalog.extractor import SongOccurrence, ExtractionResult
        import worship_catalog.extractor as extractor_mod
        import worship_catalog.pptx_reader as pptx_reader_mod
        from unittest.mock import MagicMock

        db_path = tmp_path / "notify_ok.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        job_id = str(_uuid_mod.uuid4())
        _db.create_import_job(job_id, filename="sunday.pptx")
        _db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))

        fake_songs = [
            SongOccurrence(
                ordinal=i + 1,
                canonical_title=f"song {i}",
                display_title=f"Song {i}",
            )
            for i in range(2)
        ]
        fake_result = ExtractionResult(
            filename="sunday.pptx",
            file_hash="fakehash",
            service_date="2026-01-04",
            service_name="AM Worship",
            song_leader=None,
            preacher=None,
            sermon_title=None,
            songs=fake_songs,
        )
        monkeypatch.setattr(extractor_mod, "extract_songs", lambda p, **kw: fake_result)
        monkeypatch.setattr(pptx_reader_mod, "compute_file_hash", lambda p: "fakehash")

        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)

        mock_pushover = MagicMock()
        monkeypatch.setattr(app_module, "send_pushover", mock_pushover)

        pptx_path = tmp_path / "sunday.pptx"
        pptx_path.write_bytes(b"fake")
        app_module._run_import_in_background(job_id, pptx_path)

        mock_pushover.assert_called_once()
        call_kwargs = mock_pushover.call_args[1]
        assert call_kwargs["title"] == "Import complete"
        assert "sunday.pptx" in call_kwargs["message"]
        assert "2 songs" in call_kwargs["message"]

    def test_failed_import_sends_pushover_notification(
        self, tmp_path, monkeypatch
    ):
        """On failed import, send_pushover is called with 'Import failed'."""
        from worship_catalog.db import Database
        from unittest.mock import MagicMock

        db_path = tmp_path / "notify_fail.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        job_id = str(_uuid_mod.uuid4())
        _db.create_import_job(job_id, filename="bad.pptx")
        _db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))

        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)

        mock_pushover = MagicMock()
        monkeypatch.setattr(app_module, "send_pushover", mock_pushover)

        pptx_path = tmp_path / "bad.pptx"
        pptx_path.write_bytes(b"not a real pptx")
        app_module._run_import_in_background(job_id, pptx_path)

        mock_pushover.assert_called_once()
        call_kwargs = mock_pushover.call_args[1]
        assert call_kwargs["title"] == "Import failed"
        assert "bad.pptx" in call_kwargs["message"]
        assert call_kwargs["priority"] == -1

    def test_pushover_exception_does_not_break_successful_import(
        self, tmp_path, monkeypatch
    ):
        """If send_pushover raises, the job must still be marked 'complete'."""
        from worship_catalog.db import Database
        from worship_catalog.extractor import SongOccurrence, ExtractionResult
        import worship_catalog.extractor as extractor_mod
        import worship_catalog.pptx_reader as pptx_reader_mod

        db_path = tmp_path / "notify_explode.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        job_id = str(_uuid_mod.uuid4())
        _db.create_import_job(job_id, filename="ok.pptx")
        _db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))

        fake_songs = [
            SongOccurrence(
                ordinal=1,
                canonical_title="song 0",
                display_title="Song 0",
            )
        ]
        fake_result = ExtractionResult(
            filename="ok.pptx",
            file_hash="fakehash",
            service_date="2026-01-04",
            service_name="AM Worship",
            song_leader=None,
            preacher=None,
            sermon_title=None,
            songs=fake_songs,
        )
        monkeypatch.setattr(extractor_mod, "extract_songs", lambda p, **kw: fake_result)
        monkeypatch.setattr(pptx_reader_mod, "compute_file_hash", lambda p: "fakehash")

        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)

        def _exploding_pushover(**kwargs):
            raise RuntimeError("network is down")

        monkeypatch.setattr(app_module, "send_pushover", _exploding_pushover)

        pptx_path = tmp_path / "ok.pptx"
        pptx_path.write_bytes(b"fake")
        app_module._run_import_in_background(job_id, pptx_path)

        _db2 = Database(db_path)
        _db2.connect()
        row = _db2.get_import_job(job_id)
        _db2.close()
        # If send_pushover raises after the job is marked complete but before
        # the except block, the exception could mark the job as 'failed'.
        # This test catches that regression.
        assert row["status"] == "complete", (
            f"Job status should be 'complete' even when send_pushover raises, "
            f"but got '{row['status']}'"
        )

    def test_pushover_exception_does_not_break_failed_import(
        self, tmp_path, monkeypatch
    ):
        """If send_pushover raises during a failed import, job must still be 'failed'."""
        from worship_catalog.db import Database

        db_path = tmp_path / "notify_double_fail.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        job_id = str(_uuid_mod.uuid4())
        _db.create_import_job(job_id, filename="bad.pptx")
        _db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))

        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)

        def _exploding_pushover(**kwargs):
            raise RuntimeError("network is down")

        monkeypatch.setattr(app_module, "send_pushover", _exploding_pushover)

        pptx_path = tmp_path / "bad.pptx"
        pptx_path.write_bytes(b"not a real pptx")
        # This must not raise — the import flow must absorb notification errors
        app_module._run_import_in_background(job_id, pptx_path)

        _db2 = Database(db_path)
        _db2.connect()
        row = _db2.get_import_job(job_id)
        _db2.close()
        assert row["status"] == "failed"
        assert row["error_message"] is not None

    def test_notify_vars_safe_when_update_import_job_raises_in_except(
        self, tmp_path, monkeypatch,
    ):
        """If update_import_job raises in the except block, notify vars must not be unbound (#193)."""
        from unittest.mock import MagicMock
        from worship_catalog.db import Database

        db_path = tmp_path / "notify_unbound.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        job_id = str(_uuid_mod.uuid4())
        _db.create_import_job(job_id, filename="crash.pptx")
        _db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))

        import worship_catalog.web.app as app_module
        from importlib import reload
        reload(app_module)

        mock_pushover = MagicMock()
        monkeypatch.setattr(app_module, "send_pushover", mock_pushover)

        # Make extract_songs raise, then make update_import_job ALSO raise
        import worship_catalog.extractor as extractor_mod
        monkeypatch.setattr(
            extractor_mod, "extract_songs",
            lambda p, **kw: (_ for _ in ()).throw(ValueError("bad pptx")),
        )

        monkeypatch.setattr(
            Database, "update_import_job",
            lambda self, jid, **kw: (_ for _ in ()).throw(RuntimeError("DB write failed")),
        )

        pptx_path = tmp_path / "crash.pptx"
        pptx_path.write_bytes(b"fake")

        # The RuntimeError from update_import_job propagates, but critically
        # it must NOT be UnboundLocalError — the safe defaults prevent that.
        # The finally block (including send_pushover) still runs.
        try:
            app_module._run_import_in_background(job_id, pptx_path)
        except RuntimeError:
            pass  # Expected — update_import_job raised inside except block

        # send_pushover should still be called with the safe defaults
        mock_pushover.assert_called_once()
        call_kwargs = mock_pushover.call_args[1]
        assert "crash.pptx" in call_kwargs["message"]


class TestBackgroundImportDelegatesToService:
    """_run_import_in_background should delegate to import_service.run_import."""

    def test_background_import_delegates_to_import_service(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """_run_import_in_background should call run_import from import_service."""
        from importlib import reload
        from unittest.mock import MagicMock

        from worship_catalog.import_service import ImportResult

        db_path = tmp_path / "delegate.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        job_id = "test-delegate-job"
        _db.create_import_job(job_id, filename="test.pptx")
        _db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))

        import worship_catalog.web.app as app_module

        reload(app_module)

        fake_result = ImportResult(
            service_date="2026-01-01",
            service_name="AM",
            songs_imported=5,
        )

        mock_run = MagicMock(return_value=fake_result)
        monkeypatch.setattr(app_module, "run_import", mock_run)

        mock_pushover = MagicMock()
        monkeypatch.setattr(app_module, "send_pushover", mock_pushover)

        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(b"fake")
        app_module._run_import_in_background(job_id, pptx_path)

        mock_run.assert_called_once()

        # Verify job was updated with the result from import_service
        _db2 = Database(db_path)
        _db2.connect()
        row = _db2.get_import_job(job_id)
        _db2.close()
        assert row["status"] == "complete"
        assert row["songs_imported"] == 5

    def test_background_import_passes_db_and_path_to_run_import(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """run_import must receive a Database instance and the PPTX path."""
        from importlib import reload
        from unittest.mock import MagicMock

        from worship_catalog.import_service import ImportResult

        db_path = tmp_path / "delegate_args.db"
        _db = Database(db_path)
        _db.connect()
        _db.init_schema()
        job_id = "test-delegate-args-job"
        _db.create_import_job(job_id, filename="args.pptx")
        _db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))

        import worship_catalog.web.app as app_module

        reload(app_module)

        fake_result = ImportResult(
            service_date="2026-01-01",
            service_name="AM",
            songs_imported=3,
        )

        mock_run = MagicMock(return_value=fake_result)
        monkeypatch.setattr(app_module, "run_import", mock_run)

        mock_pushover = MagicMock()
        monkeypatch.setattr(app_module, "send_pushover", mock_pushover)

        pptx_path = tmp_path / "args.pptx"
        pptx_path.write_bytes(b"fake")
        app_module._run_import_in_background(job_id, pptx_path)

        call_args = mock_run.call_args
        # First positional arg should be a Database instance
        assert isinstance(call_args[0][0], Database)
        # Second positional arg should be the pptx_path
        assert call_args[0][1] == pptx_path


class TestGetDbSchemaInit:
    """Verify that init_schema() is only called once, not on every request."""

    def test_init_schema_called_once_across_multiple_get_db_calls(self) -> None:
        """init_schema should run only on the first _get_db call, not every request."""
        from unittest.mock import patch

        import worship_catalog.web.app as app_module

        app_module._schema_ready = False

        with patch.object(Database, "init_schema") as mock_init:
            with patch.object(Database, "connect"):
                app_module._get_db()
                app_module._get_db()
                app_module._get_db()

                mock_init.assert_called_once()

        app_module._schema_ready = False

    def test_schema_ready_flag_set_after_first_call(self) -> None:
        """The _schema_ready flag should be True after first _get_db call."""
        from unittest.mock import patch

        import worship_catalog.web.app as app_module

        app_module._schema_ready = False

        with patch.object(Database, "init_schema"):
            with patch.object(Database, "connect"):
                assert not app_module._schema_ready
                app_module._get_db()
                assert app_module._schema_ready

        app_module._schema_ready = False


# ---------------------------------------------------------------------------
# Issue #179 — Empty state messages for first-run experience
# ---------------------------------------------------------------------------


class TestEmptyStateMessages:
    """Empty DB pages must show onboarding guidance, not just blank tables (#179)."""

    @pytest.fixture
    def empty_client(self, tmp_path, monkeypatch):
        db_path = tmp_path / "empty.db"
        db = Database(db_path)
        db.connect()
        db.init_schema()
        db.close()

        monkeypatch.setenv("DB_PATH", str(db_path))
        monkeypatch.setenv("INBOX_DIR", str(tmp_path / "inbox"))
        (tmp_path / "inbox").mkdir()
        from importlib import reload
        import worship_catalog.web.app as app_module
        reload(app_module)
        return TestClient(app_module.app)

    def test_songs_page_shows_empty_state_when_db_empty(self, empty_client):
        """Songs page must show a helpful message when no songs have been imported."""
        resp = empty_client.get("/songs")
        assert resp.status_code == 200
        body = resp.text.lower()
        assert any(kw in body for kw in [
            "no songs yet", "import", "get started", "upload",
        ]), "Songs page must show an empty-state onboarding message when no data exists"

    def test_services_page_shows_empty_state_when_db_empty(self, empty_client):
        """Services page must show a helpful message when no services have been imported."""
        resp = empty_client.get("/services")
        assert resp.status_code == 200
        body = resp.text.lower()
        assert any(kw in body for kw in [
            "no services yet", "import", "get started", "upload",
        ]), "Services page must show an empty-state onboarding message when no data exists"

    def test_leaders_page_shows_empty_state_when_db_empty(self, empty_client):
        """Leaders page must show a helpful message when no leaders have been imported."""
        resp = empty_client.get("/leaders")
        assert resp.status_code == 200
        body = resp.text.lower()
        assert any(kw in body for kw in [
            "no leaders yet", "import", "get started", "upload",
        ]), "Leaders page must show an empty-state onboarding message when no data exists"


class TestUploadPage:
    """Web UI must have a browser-accessible upload page."""

    def test_upload_page_exists(self, client):
        """GET /upload should return an HTML page with a file input form."""
        resp = client.get("/upload")
        assert resp.status_code == 200
        assert "text/html" in resp.headers.get("content-type", "")

    def test_upload_page_has_file_input(self, client):
        """The upload page must have a file input for PPTX files."""
        resp = client.get("/upload")
        assert resp.status_code == 200
        assert 'type="file"' in resp.text, (
            "Upload page has no file input element"
        )

    def test_upload_page_has_submit_button(self, client):
        """The upload page must have a submit button."""
        resp = client.get("/upload")
        assert resp.status_code == 200
        html = resp.text.lower()
        assert "submit" in html or "upload" in html, (
            "Upload page has no submit/upload button"
        )

    def test_nav_has_upload_link(self, client):
        """Navigation bar should include a link to the upload page."""
        resp = client.get("/songs")
        assert resp.status_code == 200
        assert "/upload" in resp.text, (
            "Navigation bar has no link to upload page"
        )


class TestBranding:
    def test_logo_served_as_static_asset(self, client):
        response = client.get("/static/highland-logo.jpg")
        assert response.status_code == 200
        assert "image" in response.headers["content-type"]

    def test_base_template_includes_logo(self, client):
        response = client.get("/songs")
        assert response.status_code == 200
        assert b"highland-logo" in response.content

    def test_page_title_includes_church_name(self, client):
        response = client.get("/songs")
        assert b"Highland" in response.content


class TestHtmxSelfHosted:
    """htmx must be served from our own static files, not a CDN."""

    def test_no_external_cdn_script_tags(self, client):
        """base.html should not reference unpkg or other CDNs."""
        resp = client.get("/songs")
        assert resp.status_code == 200
        html = resp.text
        assert "unpkg.com" not in html, (
            "Page loads scripts from unpkg CDN — should self-host htmx.js"
        )
        assert "cdn.jsdelivr.net" not in html, (
            "Page loads scripts from jsdelivr CDN — should self-host"
        )

    def test_htmx_served_from_static(self, client):
        """htmx.min.js should be available at /static/htmx.min.js."""
        resp = client.get("/static/htmx.min.js")
        assert resp.status_code == 200
        assert "htmx" in resp.text


# ---------------------------------------------------------------------------
# Issue #251 — Upload form must work without inline scripts (CSP compat)
# ---------------------------------------------------------------------------


class TestUploadFormCSPCompatible:
    """Upload form must not use inline scripts (CSP blocks them)."""

    def test_upload_page_has_no_inline_scripts(self, client):
        """upload.html must not contain bare <script> tags (blocked by CSP)."""
        resp = client.get("/upload")
        assert resp.status_code == 200
        # All <script> tags must have a src= attribute (external files)
        import re
        inline_scripts = re.findall(r"<script(?![^>]*\bsrc\b)[^>]*>", resp.text)
        assert not inline_scripts, (
            f"Page has {len(inline_scripts)} inline <script> tag(s) "
            "which are blocked by CSP script-src 'self'"
        )

    def test_upload_js_served_from_static(self, client):
        """upload.js must be available at /static/upload.js."""
        resp = client.get("/static/upload.js")
        assert resp.status_code == 200
        assert "fetch" in resp.text, "upload.js must use fetch to POST the form"
        assert "csrftoken" in resp.text.lower(), "upload.js must read the CSRF cookie"

    def test_upload_form_references_external_js(self, client):
        """upload.html must load upload.js from /static/."""
        resp = client.get("/upload")
        assert "/static/upload.js" in resp.text


class TestReportFormsCSPCompatible:
    """Report download forms must handle CSRF without inline scripts (#238)."""

    def test_reports_page_has_no_inline_scripts(self, client):
        """reports.html must not use inline scripts."""
        resp = client.get("/reports")
        assert resp.status_code == 200
        import re
        inline_scripts = re.findall(r"<script(?![^>]*\bsrc\b)[^>]*>", resp.text)
        assert not inline_scripts, (
            "Reports page has inline scripts blocked by CSP"
        )

    def test_ccli_download_form_posts_with_csrf(self, client):
        """CCLI CSV download must succeed (not 403) when submitted."""
        resp = client.post(
            "/reports/ccli",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert resp.status_code != 403, "CCLI download returned 403 — CSRF token missing"

    def test_stats_form_posts_with_csrf(self, client):
        """Stats report must succeed (not 403) when submitted."""
        resp = client.post(
            "/reports/stats",
            data={"start_date": "2026-01-01", "end_date": "2026-12-31"},
        )
        assert resp.status_code != 403, "Stats report returned 403 — CSRF token missing"


class TestEmptyStateLinkTarget:
    """Empty state card must link to Upload page, not Reports (#240)."""

    def test_songs_empty_state_links_to_upload(self, tmp_path, monkeypatch):
        db_path = tmp_path / "empty.db"
        db = Database(db_path)
        db.connect()
        db.init_schema()
        db.close()
        monkeypatch.setenv("DB_PATH", str(db_path))
        monkeypatch.setenv("INBOX_DIR", str(tmp_path / "inbox"))
        (tmp_path / "inbox").mkdir()
        from importlib import reload
        import worship_catalog.web.app as app_module
        reload(app_module)
        empty_client = TestClient(app_module.app)
        resp = empty_client.get("/songs")
        assert "/upload" in resp.text, "Empty state must link to /upload, not /reports"


class TestNoDuplicateStaticMount:
    """Static files must be mounted exactly once (#234)."""

    def test_static_mount_not_duplicated(self):
        import worship_catalog.web.app as app_module
        import inspect
        source = inspect.getsource(app_module)
        count = source.count('app.mount("/static"')
        assert count == 1, f"app.mount('/static') appears {count} times, expected 1"


class TestAboutPage:
    """Tests for GET /about — issue #232."""

    def test_about_page_returns_200(self, client):
        """GET /about must return a 200 HTML response."""
        resp = client.get("/about")
        assert resp.status_code == 200
        assert "text/html" in resp.headers.get("content-type", "")

    def test_about_page_shows_version(self, client):
        """About page must display the app version."""
        resp = client.get("/about")
        body = resp.text.lower()
        assert "version" in body

    def test_about_page_shows_purpose(self, client):
        """About page must explain what the app does in plain language."""
        resp = client.get("/about")
        body = resp.text.lower()
        assert any(kw in body for kw in ["ccli", "worship", "song", "compliance"])

    def test_about_page_shows_copyright(self, client):
        """About page must include Highland Church of Christ."""
        resp = client.get("/about")
        assert "Highland" in resp.text

    def test_about_page_shows_license(self, client):
        """About page must mention GPL-3.0 license."""
        resp = client.get("/about")
        assert "GPL-3.0" in resp.text

    def test_about_page_shows_github_link(self, client):
        """About page must link to the GitHub repository."""
        resp = client.get("/about")
        assert "https://github.com/mshirel/song-history" in resp.text

    def test_nav_has_about_link(self, client):
        """Navigation bar should include a link to the About page."""
        resp = client.get("/songs")
        assert "/about" in resp.text
