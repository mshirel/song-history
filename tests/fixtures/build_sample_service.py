"""Build the committed golden extraction fixture (#497).

Generates ``sample_service.pptx`` — a synthetic-but-realistic worship deck that
mimics the structure of a real Paperless Hymnal service export: a metadata table
slide (date / service / song leader / preacher / sermon title) followed by three
songs whose title slides carry publisher footers and credit lines.

The deck is committed (so extraction is regression-tested in CI without an
uncommittable real service file), and this script is committed too so the fixture
is reproducible. Regenerate with::

    python tests/fixtures/build_sample_service.py

The three songs deliberately exercise different credit shapes so the golden test
covers every field:

  1. "Majesty"            — "Words by:" + "Arr.:" + publisher (words_by, arranger)
  2. "Mighty To Save"     — no publisher, no credits; grouped via section prefixes
                            (all credit fields None, publisher None)
  3. "How Great Thou Art" — "Words and Music by:" + publisher (words_by == music_by)

If you change this deck, regenerate the expected JSON from the extractor and
re-commit both files (see the module docstring in test_extraction_integration.py).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# Metadata for the title/table slide. Kept in one place so it can be reused when
# regenerating the expected JSON.
SERVICE_METADATA: list[tuple[str, str]] = [
    ("Date", "2026-06-14"),
    ("Service", "Morning Worship"),
    ("Song Leader", "Matt Shirel"),
    ("Preacher", "David Morris"),
    ("Sermon Title", "The Greatness of God"),
]

# Each song is a list of slides; each slide is a list of text lines. The first
# slide of every song must be recognizable as a song-title slide — either via a
# publisher marker ("PaperlessHymnal.com") or a section prefix ("1 – Title").
SONG_SLIDES: list[list[list[str]]] = [
    # Song 1 — Majesty (publisher + Words by + Arr.)
    [
        [
            "Majesty",
            "Worship His majesty",
            "Words by: Jack Hayford",
            "Arr.: Ken Young",
            "PaperlessHymnal.com",
        ],
        [
            "Majesty",
            "Kingdom authority",
            "Flow from His throne unto His own",
        ],
    ],
    # Song 2 — Mighty To Save (no publisher, no credits; section-prefix grouped)
    [
        [
            "1 – Mighty To Save",
            "Everyone needs compassion",
            "Love that's never failing",
        ],
        [
            "2 – Mighty To Save",
            "Savior, He can move the mountains",
            "My God is mighty to save",
        ],
        [
            "3 – Mighty To Save",
            "So take me as You find me",
            "All my fears and failures",
        ],
    ],
    # Song 3 — How Great Thou Art (publisher + combined Words and Music by)
    [
        [
            "How Great Thou Art",
            "O Lord my God, when I in awesome wonder",
            "Words and Music by: Stuart K. Hine",
            "PaperlessHymnal.com",
        ],
        [
            "How Great Thou Art",
            "Then sings my soul, my Savior God to Thee",
        ],
    ],
]

OUTPUT_PATH = Path(__file__).parent / "sample_service.pptx"


def _add_text_box(slide, lines: list[str]) -> None:
    """Add a text box with one paragraph per line."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            tf.add_paragraph().text = line


def build(output_path: Path = OUTPUT_PATH) -> Path:
    """Build the sample-service deck and save it to *output_path*."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank

    # Slide 0 — metadata table (key/value rows the extractor recognizes).
    meta_slide = prs.slides.add_slide(blank_layout)
    rows = len(SERVICE_METADATA)
    table = meta_slide.shapes.add_table(
        rows, 2, Inches(1), Inches(1), Inches(8), Inches(0.4 * rows)
    ).table
    for r, (key, val) in enumerate(SERVICE_METADATA):
        table.cell(r, 0).text = key
        table.cell(r, 1).text = val

    # Song slides.
    for song in SONG_SLIDES:
        for slide_lines in song:
            slide = prs.slides.add_slide(blank_layout)
            _add_text_box(slide, slide_lines)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
