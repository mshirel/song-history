"""Build the committed self-titled non-song regression fixture (#527).

The middle unit mirrors the empty Paperless Hymnal placeholder found in the
Wesley 2026-04-05 service: every slide contains only its own title, once behind
a section prefix and once bare.  It has no lyrics, credits, or publisher footer.
Distinct real songs on either side guard against dropping or merging neighbours.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

SERVICE_METADATA: list[tuple[str, str]] = [
    ("Date", "2026-04-05"),
    ("Service", "Morning Worship"),
    ("Song Leader", "Synthetic Leader"),
    ("Preacher", "Synthetic Preacher"),
    ("Sermon Title", "Synthetic Sermon"),
]

SLIDE_LINES: list[list[str]] = [
    [
        "1 - Amazing Grace",
        "A real lyric line",
        "Words and Music by: Public Domain",
        "PaperlessHymnal.com",
    ],
    [
        "2 - Amazing Grace",
        "Another real lyric line",
        "PaperlessHymnal.com",
    ],
    ["1-1 What the Lord Has Done", "What the Lord Has Done"],
    ["1-2 What the Lord Has Done", "What the Lord Has Done"],
    ["C-1 What the Lord Has Done", "What the Lord Has Done"],
    [
        "1 - Doxology",
        "A distinct real lyric line",
        "Words and Music by: Public Domain",
        "PaperlessHymnal.com",
    ],
]

OUTPUT_PATH = Path(__file__).parent / "self_titled_nonsong_service.pptx"


def _add_text_box(slide, lines: list[str]) -> None:
    """Add one text box with one paragraph per line."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    for index, line in enumerate(lines):
        if index == 0:
            text_frame.paragraphs[0].text = line
        else:
            text_frame.add_paragraph().text = line


def build(output_path: Path = OUTPUT_PATH) -> Path:
    """Build and save the synthetic service deck."""
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]

    metadata_slide = presentation.slides.add_slide(blank_layout)
    table = metadata_slide.shapes.add_table(
        len(SERVICE_METADATA),
        2,
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(0.4 * len(SERVICE_METADATA)),
    ).table
    for row, (key, value) in enumerate(SERVICE_METADATA):
        table.cell(row, 0).text = key
        table.cell(row, 1).text = value

    for lines in SLIDE_LINES:
        slide = presentation.slides.add_slide(blank_layout)
        _add_text_box(slide, lines)

    presentation.save(output_path)
    return output_path


if __name__ == "__main__":
    print(f"Wrote {build()}")
