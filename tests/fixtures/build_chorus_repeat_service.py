"""Build the committed phantom-chorus regression fixture (#471).

Generates ``chorus_repeat_service.pptx`` — a synthetic-but-realistic Paperless
Hymnal service export that reproduces the Wesley 2026-04-05 phantom-song bug: a
verse/chorus/verse hymn whose chorus slide leads with a repeated lyric line
("He arose!  He arose!") that out-scores the real song title, so the grouping
step used to split the chorus into its own spurious one-slide "song" sandwiched
between the two verse halves.

The deck deliberately follows the phantom hymn with a genuinely DISTINCT second
song ("I Know That My Redeemer Lives") so the golden test also proves the fix
does not over-merge real back-to-back songs.

The deck and its expected JSON are both committed so extraction is regression-
tested in CI without an uncommittable real service file. Regenerate with::

    python tests/fixtures/build_chorus_repeat_service.py

If you change this deck, regenerate the expected JSON from the extractor and
re-commit both files (see the module docstring in test_extraction_integration.py).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# Metadata for the title/table slide (mirrors the real AM_Worship_2026.4.5 deck).
SERVICE_METADATA: list[tuple[str, str]] = [
    ("Date", "2026-04-05"),
    ("Service", "Morning Worship"),
    ("Song Leader", "Wesley McGehee"),
    ("Preacher", "David Morris"),
    ("Sermon Title", "The Empty Tomb"),
]

# Each song is a list of slides; each slide is a list of text lines. Slide text is
# copied from the shape of real Paperless Hymnal exports: verse slides carry a
# "N - Title" section line + a "Words and Music by:" credit + the publisher footer;
# chorus slides carry a "c - Title" section line + the publisher footer but LEAD
# with repeated lyric lines that used to win title selection (the phantom bug).
SONG_SLIDES: list[list[list[str]]] = [
    # Song 1 — Low in the Grave He Lay (verse / chorus / verse). The chorus slide is
    # the phantom trigger: its "He arose!  He arose!" lead-line is shorter than the
    # real title, so pre-fix it split into a one-slide "He arose! He arose!" song.
    [
        [
            "1.",
            "Low in the grave He lay",
            "Jesus, my Savior!",
            "Waiting the coming day",
            "Jesus my Lord!",
            "1 - Low in the Grave He Lay",
            "Words and Music by: Robert Lowry",
            "PaperlessHymnal.com",
        ],
        [
            "Up from the grave He arose",
            "With a mighty triumph o'er His foes;",
            "He arose a Victor from the dark domain,",
            "And He lives forever with His saints to reign;",
            "He arose!  He arose!",
            "Hallelujah!  Christ arose!",
            "c - Low in the Grave He Lay",
            "PaperlessHymnal.com",
        ],
        [
            "2.",
            "Vainly they watch His bed",
            "Jesus, my Savior!",
            "Vainly they seal the dead",
            "Jesus my Lord!",
            "2 - Low in the Grave He Lay",
            "Words and Music by: Robert Lowry",
            "PaperlessHymnal.com",
        ],
    ],
    # Song 2 — I Know That My Redeemer Lives (distinct; must stay its own song). A
    # verse followed by a chorus slide, proving the fix does not fold a legitimate
    # back-to-back song into the previous one.
    [
        [
            "1.",
            "I know that my Redeemer lives,",
            "And ever prays for me;",
            "I know eternal life He gives,",
            "From sin and sorrow free.",
            "1 - I Know That My Redeemer Lives",
            "Words and Music by: Fred A. Fillmore",
            "PaperlessHymnal.com",
        ],
        [
            "I know, I know that my Redeemer lives,",
            "I know, I know eternal life He gives;",
            "I know, I know that my Redeemer lives.",
            "c - I Know That My Redeemer Lives",
            "PaperlessHymnal.com",
        ],
    ],
]

OUTPUT_PATH = Path(__file__).parent / "chorus_repeat_service.pptx"


def _add_text_box(slide, lines: list[str]) -> None:
    """Add a text box with one paragraph per line."""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            tf.paragraphs[0].text = line
        else:
            tf.add_paragraph().text = line


def build(output_path: Path = OUTPUT_PATH) -> Path:
    """Build the chorus-repeat deck and save it to *output_path*."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank

    # Slide 0 — metadata table (key/value rows the extractor recognizes).
    meta_slide = prs.slides.add_slide(blank_layout)
    rows = len(SERVICE_METADATA)
    table = meta_slide.shapes.add_table(
        rows, 2, Inches(1), Inches(1), Inches(8), Inches(0.4 * rows)
    ).table
    for r, (key, val) in enumerate(SERVICE_METADATA):
        table.cell(r, 0).text = key
        table.cell(r, 1).text = val

    # Song slides.
    for song in SONG_SLIDES:
        for slide_lines in song:
            slide = prs.slides.add_slide(blank_layout)
            _add_text_box(slide, slide_lines)

    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
