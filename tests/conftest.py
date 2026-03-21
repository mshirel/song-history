"""Shared pytest fixtures for worship-catalog tests."""

import os
import socket
from collections.abc import Generator
from typing import Any

import pytest

from worship_catalog.db import Database
from worship_catalog.pptx_reader import Slide, SlideImage, SlideText

# E2E server URL — configurable via env var for CI (default: local dev server)
E2E_BASE_URL = os.environ.get("E2E_BASE_URL", "http://localhost:8000")


def _e2e_server_is_running() -> bool:
    """Return True if the E2E server at E2E_BASE_URL is accepting connections."""
    from urllib.parse import urlparse

    parsed = urlparse(E2E_BASE_URL)
    host = parsed.hostname or "localhost"
    port = parsed.port or 8000
    try:
        with socket.create_connection((host, port), timeout=2):
            return True
    except (ConnectionRefusedError, OSError):
        return False


_e2e_server_available = _e2e_server_is_running()


@pytest.fixture(scope="module")
def browser_page() -> Generator[Any, None, None]:
    """Launch a Chromium browser and yield a page. Skip if server not running."""
    if not _e2e_server_available:
        pytest.skip(
            f"No server running at {E2E_BASE_URL} — start the server to run E2E tests"
        )

    from playwright.sync_api import sync_playwright

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page()
        yield page
        browser.close()


class CsrfAwareClient:
    """Wraps TestClient to automatically include the CSRF token on POST requests."""

    def __init__(self, inner):
        self._inner = inner
        self._csrf_token = None

    def _ensure_token(self):
        if self._csrf_token is None:
            resp = self._inner.get("/songs")
            self._csrf_token = resp.cookies.get("csrftoken", "")
        return self._csrf_token or ""

    def get(self, *args, **kwargs):
        return self._inner.get(*args, **kwargs)

    def post(self, *args, **kwargs):
        token = self._ensure_token()
        headers = dict(kwargs.pop("headers", {}) or {})
        headers.setdefault("X-CSRFToken", token)
        return self._inner.post(*args, headers=headers, **kwargs)

    def __getattr__(self, name):
        return getattr(self._inner, name)


@pytest.fixture(scope="session")
def synthetic_pptx(tmp_path_factory):
    """
    Create a minimal synthetic PPTX that mimics a Paperless Hymnal worship service.

    Slide layout:
      0: metadata table slide (Service Data → alternating key/value pairs)
      1-2: Amazing Grace slides (title + verse text)
      3: How Great Thou Art slide

    This fixture exists so integration tests run in CI without requiring the real
    worship PPTX file from the data/ directory (#85).
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank layout

    def add_text_box(slide, lines, left=Inches(0.5), top=Inches(0.5),
                     width=Inches(9), height=Inches(2)):
        """Add a text box with one paragraph per line."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
        return tf

    def add_table_slide(prs, rows_data):
        """Add a slide with a table containing the given rows."""
        slide = prs.slides.add_slide(blank_layout)
        cols = 2
        rows = len(rows_data)
        table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1),
                                        Inches(8), Inches(0.4 * rows)).table
        for r, (key, val) in enumerate(rows_data):
            table.cell(r, 0).text = key
            table.cell(r, 1).text = val
        return slide

    # Slide 0: metadata table
    add_table_slide(prs, [
        ("Date", "2026-02-15"),
        ("Service", "Morning Worship"),
        ("Song Leader", "Matt"),
        ("Preacher", "Pastor John"),
        ("Sermon Title", "Grace Abounding"),
    ])

    # Slides 1-2: Amazing Grace (Paperless Hymnal format)
    # Publisher marker is required by _is_song_title_slide() to start a new song group.
    for verse_lines in [
        ["Amazing Grace", "How sweet the sound", "Words: John Newton", "PaperlessHymnal.com"],
        ["Amazing Grace", "That saved a wretch like me"],
    ]:
        slide = prs.slides.add_slide(blank_layout)
        add_text_box(slide, verse_lines)

    # Slides 3-4: How Great Thou Art
    for verse_lines in [
        ["How Great Thou Art", "O Lord my God", "Words: Stuart K. Hine", "PaperlessHymnal.com"],
        ["How Great Thou Art", "When I in awesome wonder"],
    ]:
        slide = prs.slides.add_slide(blank_layout)
        add_text_box(slide, verse_lines)

    out = tmp_path_factory.mktemp("pptx") / "AM Worship 2026.02.15.pptx"
    prs.save(out)
    return out


@pytest.fixture
def db_with_songs(tmp_path):
    """Create a minimal test DB with one service, two songs, and copy events."""
    db_path = tmp_path / "test.db"
    db = Database(db_path)
    db.connect()
    db.init_schema()

    song_id1 = db.insert_or_get_song("amazing grace", "Amazing Grace")
    song_id2 = db.insert_or_get_song("how great thou art", "How Great Thou Art")

    edition_id1 = db.insert_or_get_song_edition(
        song_id1, words_by="John Newton", music_by=None, arranger=None
    )
    edition_id2 = db.insert_or_get_song_edition(
        song_id2, words_by="Stuart K. Hine", music_by="Stuart K. Hine", arranger=None
    )

    service_id = db.insert_or_update_service(
        service_date="2026-02-15",
        service_name="AM Worship",
        source_file="test.pptx",
        source_hash="abc123",
        song_leader="Matt",
    )

    db.insert_service_song(service_id, song_id1, ordinal=1, song_edition_id=edition_id1)
    db.insert_service_song(service_id, song_id2, ordinal=2, song_edition_id=edition_id2)

    db.insert_or_get_copy_event(service_id, song_id1, "projection", song_edition_id=edition_id1)
    db.insert_or_get_copy_event(service_id, song_id2, "projection", song_edition_id=edition_id2)

    db.close()
    return db_path


@pytest.fixture
def db(db_with_songs):
    """Connected Database object pointing at the same test DB used by the client fixture."""
    database = Database(db_with_songs)
    database.connect()
    database.init_schema()
    yield database
    database.close()


def make_slide(
    index: int = 0,
    lines: list[str] | None = None,
    hidden: bool = False,
    image_blobs: list[bytes | None] | None = None,
) -> Slide:
    """Factory for creating test Slide objects without a real PPTX file."""
    text = SlideText(text_lines=lines or [])
    images = []
    if image_blobs:
        for i, blob in enumerate(image_blobs):
            images.append(SlideImage(shape_id=i + 1, blob=blob))
    return Slide(index=index, hidden=hidden, text=text, images=images)
