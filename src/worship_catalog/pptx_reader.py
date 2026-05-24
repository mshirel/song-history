"""PPTX file reading and slide parsing."""

import hashlib
import logging
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation

logger = logging.getLogger(__name__)

_HASH_CHUNK_SIZE: int = 4096
# python-pptx shape type integer for Picture shapes (MSO_SHAPE_TYPE.PICTURE == 13)
_PPTX_PICTURE_SHAPE_TYPE: int = 13


@dataclass
class SlideText:
    """Text extracted from a slide."""
    text_lines: list[str]
    """All text lines from text frames and table cells."""


@dataclass
class SlideImage:
    """Image/picture information."""
    shape_id: int
    """Unique identifier for the shape."""
    blob: bytes | None = None
    """Raw image bytes, if extracted."""


@dataclass
class Slide:
    """Parsed slide data."""
    index: int
    """0-based slide index in presentation."""
    hidden: bool
    """Whether slide is marked as hidden."""
    text: SlideText
    """Extracted text content."""
    images: list[SlideImage]
    """Picture/image shapes on slide."""


@dataclass
class ServiceMetadata:
    """Service metadata from presentation."""
    date: str | None
    """ISO format date (YYYY-MM-DD)."""
    service_name: str | None
    """Service name (e.g., 'Morning Worship')."""
    song_leader: str | None
    """Song leader name."""
    preacher: str | None
    """Preacher name."""
    sermon_title: str | None
    """Sermon title."""


def compute_file_hash(file_path: Path | str) -> str:
    """Compute SHA256 hash of file."""
    file_path = Path(file_path)
    sha256_hash = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(_HASH_CHUNK_SIZE), b""):
            sha256_hash.update(chunk)
    return sha256_hash.hexdigest()


def load_pptx(file_path: Path | str) -> Presentation:  # type: ignore[valid-type]
    """Load a PPTX file."""
    return Presentation(file_path)  # type: ignore[arg-type]


def extract_metadata_from_file(file_path: Path | str) -> dict[str, str]:
    """Extract core metadata from PPTX file path and structure."""
    file_path = Path(file_path)
    return {
        "filename": file_path.name,
        "file_hash": compute_file_hash(file_path),
    }


def is_slide_hidden(slide: Any) -> bool:
    """Check if slide has hidden attribute set."""
    try:
        # Access the slide's XML to check the show attribute
        slide_elem = slide.element
        return bool(slide_elem.get("show") == "0")
    except (AttributeError, KeyError):
        return False


def extract_text_from_slide(slide: Any) -> SlideText:
    """Extract all text from a slide's text frames and tables."""
    text_lines = []

    # Extract from text frames
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    text_lines.append(text)

        # Extract from tables
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        text_lines.append(cell_text)

    return SlideText(text_lines=text_lines)


def extract_images_from_slide(slide: Any) -> list[SlideImage]:
    """Extract image/picture shape information from slide, including raw bytes."""
    images = []

    for shape in slide.shapes:
        if shape.shape_type == _PPTX_PICTURE_SHAPE_TYPE:
            try:
                blob = shape.image.blob
            except Exception:
                logger.warning(
                    "Failed to extract blob from shape %s — returning blob=None",
                    shape.shape_id,
                    exc_info=True,
                )
                blob = None
            images.append(SlideImage(shape_id=shape.shape_id, blob=blob))

    return images


def parse_slide(slide: Any, index: int) -> Slide:
    """Parse a single slide into structured format."""
    return Slide(
        index=index,
        hidden=is_slide_hidden(slide),
        text=extract_text_from_slide(slide),
        images=extract_images_from_slide(slide),
    )


def parse_all_slides(prs: Presentation) -> list[Slide]:  # type: ignore[valid-type]
    """Parse all slides from presentation."""
    slides = []
    for i, slide in enumerate(prs.slides):  # type: ignore[attr-defined]
        slides.append(parse_slide(slide, i))
    return slides


# Date components separated by any mix of '.', '-', '/'. Used to canonicalize
# service dates to ISO YYYY-MM-DD (#387).
_DATE_PARTS_RE = re.compile(r"^(\d{4})[.\-/](\d{1,2})[.\-/](\d{1,2})$")


def normalize_service_date(raw: str | None) -> str | None:
    """Canonicalize a service date string to ISO ``YYYY-MM-DD``.

    Accepts any mix of ``.``, ``-`` and ``/`` separators and single-digit
    month/day (e.g. ``2026.5.10`` → ``2026-05-10``). Returns None for empty
    input, and the stripped original for anything that isn't a recognizable
    year-month-day triple (so we never silently drop an unexpected value).
    """
    if raw is None:
        return None
    stripped = raw.strip()
    if not stripped:
        return None
    match = _DATE_PARTS_RE.match(stripped)
    if not match:
        return stripped
    year, month, day = match.groups()
    return f"{year}-{int(month):02d}-{int(day):02d}"


# Metadata field keys (lowercased) → ServiceMetadata attribute name.
_METADATA_FIELD_KEYS: dict[str, str] = {
    "date": "date",
    "service": "service_name",
    "song leader": "song_leader",
    "preacher": "preacher",
    "sermon title": "sermon_title",
}
# Header/section labels that are never values.
_METADATA_HEADER_KEYS: frozenset[str] = frozenset({"service data", "metadata", "service info"})


def extract_metadata_from_table_slide(slide_text_lines: list[str]) -> ServiceMetadata:
    """
    Extract service metadata from a table slide.

    Keys and values live on separate lines. Rather than assume strict
    even/odd key-value alignment (which a stray header/blank line silently
    breaks, #385), scan for each recognized key and take the next non-empty
    line that is not itself a key/header as its value.
    """
    result = ServiceMetadata(
        date=None,
        service_name=None,
        song_leader=None,
        preacher=None,
        sermon_title=None,
    )

    n = len(slide_text_lines)
    for idx, raw in enumerate(slide_text_lines):
        key = raw.strip().lower()
        attr = _METADATA_FIELD_KEYS.get(key)
        if attr is None or getattr(result, attr) is not None:
            continue
        # Value = next non-empty line before the next key/header.
        for j in range(idx + 1, n):
            cand = slide_text_lines[j].strip()
            if not cand:
                continue
            cand_lower = cand.lower()
            if cand_lower in _METADATA_FIELD_KEYS or cand_lower in _METADATA_HEADER_KEYS:
                break  # next key reached — this field has no value
            setattr(result, attr, cand)
            break

    return result


def parse_filename_for_metadata(filename: str) -> ServiceMetadata:
    """
    Fallback: parse metadata from filename pattern like 'AM Worship 2026.02.15.pptx'.

    Expected pattern:
    - AM Worship YYYY.MM.DD.pptx → Morning Worship, date
    - PM Worship YYYY.MM.DD.pptx → Evening Worship, date
    """
    result = ServiceMetadata(
        date=None,
        service_name=None,
        song_leader=None,
        preacher=None,
        sermon_title=None,
    )

    # Pattern: (AM|PM) Worship YYYY.MM.DD — accepts underscores and job_id prefixes
    # (#265), single-digit month/day, and any of '.', '-', '/' separators (#387).
    match = re.search(
        r"([AP]M)[\s_]+Worship[\s_]+(\d{4}[.\-/]\d{1,2}[.\-/]\d{1,2})", filename
    )
    if match:
        am_pm, date_str = match.groups()
        result.service_name = "Morning Worship" if am_pm == "AM" else "Evening Worship"
        result.date = normalize_service_date(date_str)

    return result


def extract_service_metadata(
    first_slide: Slide, filename: str
) -> ServiceMetadata:
    """
    Extract service metadata from presentation.

    Strategy:
    1. Try to parse structured metadata from first slide table
    2. Fallback to filename parsing
    3. Return result with None for missing fields
    """
    # Try table-based extraction first
    metadata = ServiceMetadata(
        date=None, service_name=None, song_leader=None, preacher=None, sermon_title=None,
    )
    if first_slide.text.text_lines:
        metadata = extract_metadata_from_table_slide(first_slide.text.text_lines)

    # Fill any None fields from filename fallback (#169)
    if metadata.date is None or metadata.service_name is None:
        fallback = parse_filename_for_metadata(filename)
        if metadata.date is None:
            metadata.date = fallback.date
        if metadata.service_name is None:
            metadata.service_name = fallback.service_name

    # Canonicalize the date to ISO YYYY-MM-DD regardless of source (#387).
    metadata.date = normalize_service_date(metadata.date)

    return metadata
